from __future__ import annotations
import logging
from datetime import date, datetime, timedelta
import time
from pathlib import Path
from urllib.parse import quote_plus
import base64
import hashlib
import io
import json
import os
import random
import re
import secrets
import string
import uuid
import functools
from core.authorization import can_action as _core_can_action
from core.navigation import ROLE_NAVIGATION
from core.scope import restrict_user_frame as _restrict_user_frame, allowed_user_ids as _allowed_user_ids
from core.schema_contract import contract_report as _schema_contract_report
from core.health import application_health as _application_health
from core.observability import timed as _timed, performance_snapshot as _performance_snapshot
from core.repository import Repository as _Repository
from core.database_gateway import get_engine, exec_sql, query_sql
from core.auth_provider import SupabaseAuthProvider as _SupabaseAuthProvider
from core.security import password_errors as _password_errors, valid_email as _valid_email
from core.design_system import page_kicker as _page_kicker, role_presentation as _role_presentation
from core.production import page_execution as _page_execution, production_config_report as _production_config_report
from core.access_policy import allowed_user_ids as _policy_allowed_user_ids, filter_frame as _policy_filter_frame, ROLE_ALLOWED_SCOPES, ORG_ROLES
from core.system_write import system_write, is_system_write
from core.migrations import run_pending_migrations
import pandas as pd
import qrcode
import streamlit as st
import streamlit.components.v1 as components

try:
    from argon2 import PasswordHasher
except Exception:
    PasswordHasher = None
try:
    from supabase import create_client
except Exception:
    create_client = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None
try:
    import docx
except Exception:
    docx = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    import openpyxl
except Exception:
    openpyxl = None
from sqlalchemy import create_engine, text
from sqlalchemy.engine import Engine
APP_TITLE = 'Pakistan Shipping Bureau'
APP_SUBTITLE = 'World-Class Classification Society Training, Competency, Authorization and Workforce Platform'
DATABASE_URL = os.getenv('DATABASE_URL', 'sqlite:///psb_hrdm_world_class.db')
PUBLIC_URL = os.getenv('PUBLIC_URL', 'https://training.psbureau.org')
VERIFY_PUBLIC_URL = os.getenv('VERIFY_PUBLIC_URL', PUBLIC_URL)
SUPABASE_URL = os.getenv('SUPABASE_URL', '')
SUPABASE_SERVICE_ROLE_KEY = os.getenv('SUPABASE_SERVICE_ROLE_KEY', '')
SUPABASE_BUCKET = os.getenv('SUPABASE_BUCKET', 'psb-hrdm-files')
SUPABASE_ANON_KEY = os.getenv('SUPABASE_ANON_KEY', '')
AUTH_CALLBACK_URL = os.getenv('AUTH_CALLBACK_URL', '')
AUTH_MODE = os.getenv('AUTH_MODE', 'local')
SESSION_IDLE_MINUTES = int(os.getenv('SESSION_IDLE_MINUTES', '30'))
SESSION_MAX_HOURS = int(os.getenv('SESSION_MAX_HOURS', '12'))
ARGON2_TIME_COST = int(os.getenv('ARGON2_TIME_COST', '3'))
ARGON2_MEMORY_COST = int(os.getenv('ARGON2_MEMORY_COST', '65536'))
ARGON2_PARALLELISM = int(os.getenv('ARGON2_PARALLELISM', '2'))
PASSWORD_HASHER = PasswordHasher(time_cost=ARGON2_TIME_COST, memory_cost=ARGON2_MEMORY_COST, parallelism=ARGON2_PARALLELISM) if PasswordHasher else None
LOGO_PATH = Path('assets/psb-logo.png')
LOCAL_UPLOAD_DIR = Path('local_uploads')
INITIAL_ADMIN_NAME = os.getenv('INITIAL_ADMIN_NAME', 'PSB Administrator')
INITIAL_ADMIN_EMAIL = os.getenv('INITIAL_ADMIN_EMAIL', '')
INITIAL_ADMIN_LOGIN = os.getenv('INITIAL_ADMIN_LOGIN', '')
INITIAL_ADMIN_PASSWORD = os.getenv('INITIAL_ADMIN_PASSWORD', '')
ENABLE_DEMO_SEED = os.getenv('ENABLE_DEMO_SEED', 'false').strip().lower() == 'true'
DEMO_PASSWORD = os.getenv('DEMO_PASSWORD', '')
MAX_LOGIN_ATTEMPTS = int(os.getenv('MAX_LOGIN_ATTEMPTS', '5'))
LOGIN_BLOCK_MINUTES = int(os.getenv('LOGIN_BLOCK_MINUTES', '5'))
SAFE_TABLE_RE = re.compile('^[a-zA-Z0-9_]+$')
SAFE_FILENAME_RE = re.compile('[^a-zA-Z0-9._-]+')
SUPABASE_AUTH_PROVIDER = _SupabaseAuthProvider(SUPABASE_URL, SUPABASE_ANON_KEY, create_client)

class RateLimiter:

    def __init__(self):
        self._calls: dict[str, dict[str, list[float]]] = {}

    def _prune(self, lst: list[float], window: float) -> list[float]:
        cutoff = time.time() - window
        return [t for t in lst if t >= cutoff]

    def allowed(self, action: str, key: str, limit: int, per_seconds: int) -> bool:
        now_ts = time.time()
        self._calls.setdefault(action, {})
        self._calls[action].setdefault(key, [])
        calls = self._prune(self._calls[action][key], per_seconds)
        if len(calls) >= limit:
            return False
        calls.append(now_ts)
        self._calls[action][key] = calls
        return True
RATE_LIMITER = RateLimiter()
APP_ENV = os.getenv('APP_ENV', 'production' if os.getenv('RENDER') else 'local').lower()

def is_render_runtime() -> bool:
    return bool(os.getenv('RENDER') or os.getenv('RENDER_SERVICE_ID') or os.getenv('RENDER_EXTERNAL_URL'))

def database_is_persistent() -> bool:
    url = DATABASE_URL.lower().strip()
    return url.startswith(('postgresql://', 'postgresql+psycopg2://', 'postgres://'))

def storage_is_persistent() -> bool:
    return bool(SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY and SUPABASE_BUCKET)

def require_persistent_backend() -> None:
    """Prevent data loss on Render by blocking temporary SQLite/local storage."""
    if is_render_runtime() and (not database_is_persistent()):
        st.error('Persistent database is not configured. Render local SQLite storage is temporary and data will disappear after restart/redeploy.')
        st.markdown('\n        **Fix in Render → Environment Variables:**\n        ```text\n        DATABASE_URL=postgresql://postgres:PASSWORD@HOST:5432/postgres\n        SUPABASE_URL=https://your-project.supabase.co\n        SUPABASE_SERVICE_ROLE_KEY=your-service-role-key\n        SUPABASE_BUCKET=psb-hrdm-files\n        ```\n        ')
        st.stop()

STANDARDS = ['IMO RO Code', 'ISO 9001', 'ISO/IEC 17020', 'IACS PR7', 'Competency-Based Qualification System']
ROLES = ['GM', 'Admin', 'Trainer', 'Department Manager', 'NSC Surveyor', 'In-Service Surveyor', 'Surveyor', 'Plan Appraiser', 'QMS Auditor', 'Industrial Surveyor', 'Rule Development Rep', 'QMR', 'Management', 'Trainee', 'On Probation']
DEPARTMENTS = ['Survey NSC', 'NSC Survey', 'Survey Inservice', 'Plan Appraisal', 'QMS', 'Rule Development', 'Training', 'Administration']
ADMIN_ACCOUNT_STATUSES = ['Invited', 'Active', 'Suspended', 'Deactivated']
AVAILABILITY_STATUSES = ['Available', 'Busy', 'On Leave', 'Unavailable']
PERMISSION_ACTIONS = ['View', 'Create', 'Edit', 'Submit', 'Assign', 'Review', 'Approve', 'Reject', 'Close', 'Export', 'Manage']
PERMISSION_SCOPES = ['Own', 'Assigned', 'Department', 'Multiple Departments', 'Organization-wide']
PERMISSION_MODULES = ['Dashboard', 'Employee Profile', 'Development Plans', 'Succession Planning', 'Workforce Planning', 'Users & Roles', 'Departments', 'Training Dashboard', 'Training', 'Training Matrix', 'CPD', 'Competency', 'Practical / Witness', 'Gap Advisor', 'NCR / Corrective Action', 'Knowledge Library', 'Authorization', 'CRB', 'Technical Authority', 'Restrictions', 'Annual Review', 'Revalidation', 'Technical Reviews', 'QMS', 'Accreditation Readiness', 'Interpretation Portal', 'Job Allocation', 'Client Feedback', 'Performance & KPI', 'Administration']
DEFAULT_ROLE_DESCRIPTIONS = {'GM': 'Executive governance and organization-wide oversight.', 'Admin': 'System administration, identity, access and organizational configuration.', 'Trainer': 'Owns assigned learner qualification paths, training, mentoring and development progression.', 'Department Manager': 'Department-scoped qualification oversight for the assigned technical department.', 'NSC Surveyor': 'New ship construction survey professional.', 'In-Service Surveyor': 'In-service ship survey professional.', 'Surveyor': 'Surveyor progressing through or holding qualification in the selected survey path.', 'Plan Appraiser': 'Plan appraisal professional progressing through or holding the Plan Appraiser path.', 'QMS Auditor': 'QMS audit and compliance activities.', 'Industrial Surveyor': 'Industrial survey professional progressing through or holding the Industrial Surveyor path.', 'Rule Development Rep': 'Rule development and interpretation activities.', 'QMR': 'Quality management oversight.', 'Management': 'Management oversight and final business approvals.', 'Trainee': 'Person progressing through a predefined qualification path.', 'On Probation': 'Probationary person progressing through a predefined qualification path.'}
TRAINEE_PATHS = ['NSC Surveyor', 'In-Service Surveyor', 'Industrial Surveyor', 'Plan Appraiser']
JOB_TYPES = ['New Building Survey', 'In-Service Survey', 'Plan Appraisal', 'Internal Audit', 'External Audit', 'Industrial Survey', 'Rule Development', 'Witness Survey']
SCOPES = ['Hull NB', 'Hull IS', 'Machinery NB', 'Machinery IS', 'Electrical NB', 'Electrical IS', 'Statutory SOLAS', 'Statutory MARPOL', 'Plan Approval Hull', 'Plan Approval Machinery', 'Plan Approval Electrical', 'Internal Auditor', 'External Auditor', 'Industrial Surveyor', 'Rule Development']
COMPETENCY_LEVELS = ['Level 0 - Trainee', 'Level 1 - Witness Eligible', 'Level 2 - Supervised Eligible', 'Level 3 - Authorized', 'Level 4 - Senior Authorized', 'Level 5 - Principal / Lead']
FILE_CATEGORIES = ['Training Material', 'SCORM Package', 'Rule Document', 'Knowledge Bulletin', 'Survey Evidence', 'Plan Review Evidence', 'Witness Evidence', 'Certificate Template', 'Issued Certificate', 'CAPA Evidence', 'Other']
ALLOWED_EXTENSIONS = ['pdf', 'ppt', 'pptx', 'txt', 'doc', 'docx', 'xls', 'xlsx', 'png', 'jpg', 'jpeg', 'mp4', 'csv']
MAX_UPLOAD_BYTES = int(os.getenv('MAX_UPLOAD_BYTES', str(25 * 1024 * 1024)))
CORE_THEORETICAL_MODULES = [('CORE-001', 'PSB Induction and Code of Ethics', 'All', 'Core', 2), ('CORE-002', 'IMO Recognized Organization Code Awareness', 'All', 'Core', 3), ('CORE-003', 'ISO 9001 Quality Management System', 'All', 'QMS', 3), ('CORE-004', 'ISO/IEC 17020 Inspection Body Requirements', 'All', 'QMS', 3), ('CORE-005', 'IACS PR7 Training and Qualification Principles', 'All', 'Competency', 2), ('CORE-006', 'Document Control and Record Retention', 'All', 'QMS', 2), ('CORE-007', 'HSE, Risk Assessment and Site Safety', 'All', 'Safety', 3), ('CORE-008', 'Survey Reporting and Deficiency Management', 'Surveyor', 'Survey', 3), ('TECH-001', 'Hull Rules and Structural Survey Principles', 'Hull Surveyor', 'Technical', 5), ('TECH-002', 'Machinery Rules and Machinery Survey Principles', 'Machinery Surveyor', 'Technical', 5), ('TECH-003', 'Electrical Rules and Electrical Survey Principles', 'Electrical Surveyor', 'Technical', 5), ('STAT-001', 'SOLAS Statutory Survey Requirements', 'Statutory Surveyor', 'Statutory', 5), ('STAT-002', 'MARPOL Pollution Prevention Requirements', 'Statutory Surveyor', 'Statutory', 4), ('PLAN-001', 'Plan Appraisal Rule Interpretation', 'Plan Appraiser', 'Plan Appraisal', 4), ('PLAN-002', 'Plan Review Commenting and Approval Workflow', 'Plan Appraiser', 'Plan Appraisal', 3), ('AUD-001', 'Internal Audit Techniques and CAPA', 'Auditor', 'Audit', 4), ('RULE-001', 'Rule Development, Technical Circulars and Change Impact', 'Rule Development Rep', 'Rule Development', 4)]
DEFAULT_AUTH_MATRIX = [('Hull NB', 'New Building Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Hull IS', 'In-Service Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Machinery NB', 'New Building Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Machinery IS', 'In-Service Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Electrical NB', 'New Building Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Electrical IS', 'In-Service Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Statutory SOLAS', 'Statutory Survey', 2, 1, 0, 0, 3, 3, 'High', 36), ('Statutory MARPOL', 'Statutory Survey', 2, 1, 0, 0, 3, 3, 'High', 36), ('Plan Approval Hull', 'Plan Appraisal', 0, 0, 2, 1, 3, 3, 'Medium', 36), ('Plan Approval Machinery', 'Plan Appraisal', 0, 0, 2, 1, 3, 3, 'Medium', 36), ('Plan Approval Electrical', 'Plan Appraisal', 0, 0, 2, 1, 3, 3, 'Medium', 36), ('Internal Auditor', 'Internal Audit', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('External Auditor', 'External Audit', 2, 1, 0, 0, 4, 4, 'High', 36), ('Industrial Surveyor', 'Industrial Survey', 2, 1, 0, 0, 3, 3, 'Medium', 36), ('Rule Development', 'Rule Development', 1, 1, 0, 0, 4, 4, 'High', 36)]

def now() -> str:
    return datetime.now().strftime('%Y-%m-%d %H:%M:%S')

def today() -> str:
    return date.today().strftime('%Y-%m-%d')

def uid(prefix: str) -> str:
    return f'{prefix}-{uuid.uuid4().hex[:8].upper()}'

def clean(v) -> str:
    if v is None:
        return ''
    try:
        if pd.isna(v):
            return ''
    except Exception:
        pass
    return str(v)

def validate_table_name(table: str) -> None:
    if not SAFE_TABLE_RE.fullmatch(clean(table)):
        raise ValueError('Invalid table name.')

def sanitize_path_component(value: str) -> str:
    value = clean(value).replace('\\', '/').split('/')[-1]
    value = SAFE_FILENAME_RE.sub('_', value).strip('._-')
    return value or 'item'

def sanitize_filename(name: str) -> str:
    name = clean(name).replace('\\', '/').split('/')[-1]
    name = SAFE_FILENAME_RE.sub('_', name).strip('._-')
    return name or 'upload'

def validate_email(email: str) -> bool:
    return bool(re.fullmatch('[^@ \\t\\r\\n]+@[^@ \\t\\r\\n]+\\.[^@ \\t\\r\\n]+', clean(email)))

def _session_hash(token: str) -> str:
    return hashlib.sha256(clean(token).encode('utf-8')).hexdigest()







def days_until(date_text: str) -> int:
    if not clean(date_text):
        return 9999
    try:
        return (datetime.strptime(clean(date_text)[:10], '%Y-%m-%d').date() - date.today()).days
    except Exception:
        return 9999

def add_months(base_or_months, months: int | None = None) -> str:
    """Add months to today (legacy form) or to an explicit YYYY-MM-DD/date value."""
    if months is None:
        months = int(base_or_months)
        d = date.today()
    else:
        if isinstance(base_or_months, date):
            d = base_or_months
        else:
            d = datetime.strptime(str(base_or_months)[:10], '%Y-%m-%d').date()
    month = d.month - 1 + int(months)
    year = d.year + month // 12
    month = month % 12 + 1
    leap = year % 4 == 0 and (year % 100 != 0 or year % 400 == 0)
    day = min(d.day, [31, 29 if leap else 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31][month - 1])
    return date(year, month, day).strftime('%Y-%m-%d')






@st.cache_resource
def get_supabase_client():
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY or create_client is None:
        return None
    return create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)

def exec_sql(sql: str, params: dict | None=None) -> None:
    with get_engine().begin() as conn:
        conn.execute(text(sql), params or {})

def query_sql(sql: str, params: dict | None=None) -> pd.DataFrame:
    with _timed('database.query'):
        with get_engine().begin() as conn:
            return pd.read_sql(text(sql), conn, params=params or {})

def _scope_table_for_read(table: str) -> bool:
    try:
        from core.access_policy import TABLE_SCOPE_MODULES, TABLE_MUTATION_MODULES
        return table in TABLE_SCOPE_MODULES
    except Exception:
        return False

def db_all_unscoped(table: str) -> pd.DataFrame:
    validate_table_name(table)
    return REPOSITORY.select_all(table)

def db_where_unscoped(table: str, where_sql: str, params_tuple: tuple[tuple[str, object], ...]=()) -> pd.DataFrame:
    validate_table_name(table)
    return REPOSITORY.select_where(table, where_sql, dict(params_tuple))

@st.cache_data(ttl=15, show_spinner=False)
def db_count(table: str, where_sql: str='', params_tuple: tuple[tuple[str, object], ...]=()) -> int:
    try:
        validate_table_name(table)
        return REPOSITORY.count(table, where_sql, dict(params_tuple))
    except Exception:
        return 0

@st.cache_data(ttl=20, show_spinner=False)
def db_all(table: str) -> pd.DataFrame:
    try:
        frame = db_all_unscoped(table)
        actor = st.session_state.get('user') if hasattr(st, 'session_state') else None
        if actor and _scope_table_for_read(table) and not frame.empty:
            users = db_all_unscoped('users')
            uds = db_all_unscoped('user_departments') if table != 'user_departments' else frame
            
            from core.access_policy import filter_frame
            return filter_frame(frame, actor, users, uds)
        return frame
    except Exception:
        return pd.DataFrame()

@st.cache_data(ttl=20, show_spinner=False)
def db_where(table: str, where_sql: str, params_tuple: tuple[tuple[str, object], ...]=()) -> pd.DataFrame:
    try:
        frame = db_where_unscoped(table, where_sql, params_tuple)
        actor = st.session_state.get('user') if hasattr(st, 'session_state') else None
        if actor and _scope_table_for_read(table) and not frame.empty:
            users = db_all_unscoped('users')
            uds = db_all_unscoped('user_departments') if table != 'user_departments' else frame
            
            from core.access_policy import filter_frame
            return filter_frame(frame, actor, users, uds)
        return frame
    except Exception:
        return pd.DataFrame()

def clear_db_cache() -> None:
    """Clear Streamlit data caches after write operations.
    The previous version accidentally called itself recursively, which could
    freeze the app after inserts/updates/deletes.
    """
    try:
        db_all.clear()
        db_where.clear()
        db_count.clear()
    except Exception:
        pass
REPOSITORY = _Repository(query_sql=query_sql, exec_sql=exec_sql, cache_clear=clear_db_cache)

def first_row(df: pd.DataFrame) -> dict | None:
    if df is None or df.empty:
        return None
    return df.iloc[0].to_dict()

def convert_numpy_types(row: dict) -> dict:
    """Convert numpy types to Python native types for database compatibility."""
    converted = {}
    for key, value in row.items():
        if value is None:
            converted[key] = None
        elif hasattr(value, 'item'):
            converted[key] = value.item()
        else:
            converted[key] = value
    return converted

_SERVER_INTERNAL_TABLES = {
    'schema_migrations','audit_trail','auth_sessions','login_security_state',
    'scheduler_runs','qr_verification_events','notifications'
}

def _mutation_guard(table: str, action: str, row: dict | None) -> None:
    """Mandatory mutation boundary. Business writes require record-level RBAC;
    internal system writes require an explicit system_write context. Startup before
    a user session exists remains permitted for migrations/bootstrap only.
    """
    # Explicit trusted-service contexts are the only supported way for
    # authentication and other server-owned workflows to update a business
    # table without borrowing the signed-in user's module permissions.
    if is_system_write():
        return
    actor = st.session_state.get('user') if hasattr(st, 'session_state') else None
    if table in _SERVER_INTERNAL_TABLES:
        if not isinstance(actor, dict) or not actor.get('user_id'):
            return
        raise PermissionError(f'Internal table {table} can only be modified by a trusted system service.')
    if not isinstance(actor, dict) or not actor.get('user_id'):
        return
    from core.access_policy import TABLE_MUTATION_MODULES
    from core.authorization import authorize_action
    module = TABLE_MUTATION_MODULES.get(table)
    if table == 'files' and str((row or {}).get('linked_table') or '') == 'trainings':
        module = 'Training'
    if table == 'user_assignments' and str(actor.get('role') or '') == 'Trainer':
        module = 'Training'
    if not module:
        raise PermissionError(f'No mutation policy is registered for table {table}.')
    if not authorize_action(actor, module, action, row or {}, db_all=db_all_unscoped, db_where=db_where_unscoped, actor_get=actor_get):
        raise PermissionError(f'Not authorized to {action} {module} for this record.')

def db_insert(table: str, row: dict) -> None:
    _mutation_guard(table, 'Create', row)
    REPOSITORY.insert(table, row)

def db_insert_many(table: str, rows: list[dict]) -> None:
    for row in rows:
        _mutation_guard(table, 'Create', row)
    REPOSITORY.insert_many(table, rows)

def db_update(table: str, id_col: str, id_val: str, row: dict) -> None:
    existing = None
    try:
        current = db_where(table, f'{id_col} = :id', (('id', id_val),))
        if not current.empty:
            existing = current.iloc[0].to_dict()
    except Exception:
        existing = None
    merged = dict(existing or {})
    merged.update(row or {})
    if id_col not in merged:
        merged[id_col] = id_val
    _mutation_guard(table, 'Edit', merged)
    REPOSITORY.update(table, id_col, id_val, row)

def db_delete(table: str, id_col: str, id_val: str) -> None:
    existing = None
    try:
        current = db_where(table, f'{id_col} = :id', (('id', id_val),))
        if not current.empty:
            existing = current.iloc[0].to_dict()
    except Exception:
        existing = None
    _mutation_guard(table, 'Delete', existing or {id_col: id_val})
    REPOSITORY.delete(table, id_col, id_val)

@st.cache_resource(show_spinner=False)









def _malware_scan_upload(filename: str, data: bytes) -> str:
    """Scan uploads with ClamAV when available; fail closed when production requires it."""
    import shutil, subprocess, tempfile
    required = os.getenv('REQUIRE_MALWARE_SCAN', 'true' if is_render_runtime() else 'false').strip().lower() in {'1','true','yes','on'}
    scanner = shutil.which('clamscan')
    if not scanner:
        if required:
            raise RuntimeError('Malware scanning is required but ClamAV is not installed/configured.')
        return 'Validated-NoScanner-Development'
    suffix='.'+(filename.rsplit('.',1)[-1] if '.' in filename else 'bin')
    tmp=None
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(data); tmp=f.name
        result=subprocess.run([scanner,'--no-summary',tmp],capture_output=True,text=True,timeout=30)
        if result.returncode == 1:
            raise ValueError('Upload rejected: malware signature detected.')
        if result.returncode != 0:
            if required: raise RuntimeError('Malware scanner failed; upload blocked.')
            return 'Validated-ScannerUnavailable'
        return 'Validated-Clean'
    finally:
        if tmp:
            try: Path(tmp).unlink(missing_ok=True)
            except Exception: pass

def _validate_ooxml_archive(data: bytes) -> None:
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            infos=z.infolist()
            if len(infos)>2000: raise ValueError('Office document contains too many embedded entries.')
            total=sum(int(i.file_size or 0) for i in infos)
            if total > 100*1024*1024: raise ValueError('Office document expands beyond the safe processing limit.')
    except zipfile.BadZipFile:
        raise ValueError('Invalid Office document container.')

def _validate_upload_bytes(filename: str, ext: str, mime_type: str, data: bytes) -> None:
    if len(data) > MAX_UPLOAD_BYTES:
        raise ValueError(f'File exceeds the {MAX_UPLOAD_BYTES // (1024*1024)} MB upload limit.')
    if not data:
        raise ValueError('Empty files are not allowed.')
    # Lightweight magic-byte validation. Office OOXML files are ZIP containers.
    sig = data[:12]
    checks = {
        'pdf': lambda: sig.startswith(b'%PDF-'),
        'png': lambda: sig.startswith(b'\\x89PNG\\r\\n\\x1a\\n'),
        'jpg': lambda: sig.startswith(b'\\xff\\xd8\\xff'),
        'jpeg': lambda: sig.startswith(b'\\xff\\xd8\\xff'),
        'docx': lambda: sig.startswith(b'PK'),
        'xlsx': lambda: sig.startswith(b'PK'),
        'pptx': lambda: sig.startswith(b'PK'),
    }
    if ext in checks and not checks[ext]():
        raise ValueError(f'File content does not match the .{ext} extension.')
    if ext in {'docx','xlsx','pptx'}:
        _validate_ooxml_archive(data)
    # Reject executable/script signatures and embedded HTML uploads.
    if sig.startswith((b'MZ', b'\\x7fELF')) or ext in {'exe','dll','js','html','htm','svg'}:
        raise ValueError('Executable/script content is not permitted.')

def secure_file_url(file_row: dict, expires_seconds: int = 600) -> str:
    """Return a short-lived signed URL for private Supabase objects.

    Local files intentionally do not return a browser URL; callers should use a
    download button after reading bytes server-side.
    """
    provider = clean((file_row or {}).get('storage_provider'))
    path = clean((file_row or {}).get('storage_path'))
    if provider == 'supabase' and path:
        client = get_supabase_client()
        if client is None:
            return ''
        try:
            result = client.storage.from_(SUPABASE_BUCKET).create_signed_url(path, max(60, min(int(expires_seconds), 900)))
            if isinstance(result, dict):
                return clean(result.get('signedURL') or result.get('signedUrl') or result.get('signed_url'))
        except Exception:
            return ''
    return ''

def secure_file_bytes(file_row: dict) -> bytes | None:
    provider = clean((file_row or {}).get('storage_provider'))
    path = clean((file_row or {}).get('storage_path'))
    try:
        if provider == 'local' and path:
            full = LOCAL_UPLOAD_DIR / path
            if full.exists() and full.is_file():
                return full.read_bytes()
        if provider == 'supabase' and path:
            client = get_supabase_client()
            if client is not None:
                return client.storage.from_(SUPABASE_BUCKET).download(path)
    except Exception:
        return None
    return None

def upload_file(uploaded_file, actor: dict, linked_table: str, linked_id: str, category: str, information_classification: str='Internal') -> dict:
    user_key = actor_get(actor, 'user_id') or 'anon'
    if not RATE_LIMITER.allowed('upload', user_key, limit=15, per_seconds=60):
        raise RuntimeError('Rate limit exceeded for uploads. Try again later.')
    file_id = uid('FILE')
    ext = uploaded_file.name.split('.')[-1].lower() if '.' in uploaded_file.name else ''
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f'File type .{ext} is not allowed.')
    data = uploaded_file.getvalue()
    _validate_upload_bytes(uploaded_file.name, ext, uploaded_file.type or '', data)
    security_status = _malware_scan_upload(uploaded_file.name, data)
    information_classification = clean(information_classification) or 'Internal'
    if information_classification not in {'Public','Internal','Confidential','Restricted Technical'}:
        raise ValueError('Invalid information classification.')
    linked_table = sanitize_path_component(linked_table)
    linked_id = sanitize_path_component(linked_id)
    filename = sanitize_filename(uploaded_file.name)
    storage_path = f'{sanitize_path_component(category)}/{linked_table}/{linked_id}/{file_id}_{filename}'
    provider = 'local'
    client = get_supabase_client()
    if client is not None:
        try:
            try:
                # Evidence/training storage is private. Public certificate verification
                # uses the verification service, not a public storage bucket.
                client.storage.create_bucket(SUPABASE_BUCKET, options={'public': False})
            except Exception:
                try:
                    client.storage.update_bucket(SUPABASE_BUCKET, options={'public': False})
                except Exception:
                    pass
            client.storage.from_(SUPABASE_BUCKET).upload(storage_path, data, {'content-type': uploaded_file.type or 'application/octet-stream', 'upsert': 'false'})
            provider = 'supabase'
        except Exception as e:
            if is_render_runtime():
                raise RuntimeError(f'Supabase Storage upload failed on Render. Configure SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY and a PRIVATE SUPABASE_BUCKET. Details: {e}')
            provider = 'local'
    if provider == 'local':
        if is_render_runtime():
            raise RuntimeError('Local file storage is disabled on Render because it is temporary. Configure private Supabase Storage.')
        local_path = LOCAL_UPLOAD_DIR / storage_path
        local_path.parent.mkdir(parents=True, exist_ok=True)
        local_path.write_bytes(data)
    extracted = extract_text(uploaded_file.name, data)
    row = {'file_id': file_id, 'owner_user_id': actor_get(actor, 'user_id'), 'owner_name': actor_get(actor, 'name'), 'linked_table': linked_table, 'linked_id': linked_id, 'category': category, 'file_name': filename, 'file_ext': ext, 'mime_type': uploaded_file.type or '', 'storage_provider': provider, 'storage_path': storage_path, 'public_url': '', 'extracted_text': extracted[:10000], 'ocr_status': 'Extracted' if extracted else 'Pending/Not Supported', 'review_status': 'Pending Review', 'security_status': security_status, 'information_classification': information_classification, 'size_bytes': len(data), 'created_on': now(), 'updated_on': now()}
    db_insert('files', row)
    audit('File Uploaded', f'{filename} linked to {linked_table}:{linked_id}', actor=actor)
    return row

def extract_text(name: str, data: bytes) -> str:
    lower = name.lower()
    try:
        if lower.endswith(('.txt', '.csv')):
            return data.decode('utf-8', errors='ignore')
        if lower.endswith('.pdf') and PdfReader:
            reader = PdfReader(io.BytesIO(data))
            return '\n'.join([p.extract_text() or '' for p in reader.pages])
        if lower.endswith('.docx') and docx:
            doc = docx.Document(io.BytesIO(data))
            return '\n'.join((p.text for p in doc.paragraphs))
        if lower.endswith('.pptx') and Presentation:
            prs = Presentation(io.BytesIO(data))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        texts.append(shape.text)
            return '\n'.join(texts)
        if lower.endswith(('.xlsx', '.xlsm')) and openpyxl:
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            texts=[]
            for ws in wb.worksheets:
                texts.append(f'Worksheet: {ws.title}')
                for row in ws.iter_rows(values_only=True):
                    vals=[str(v).strip() for v in row if v is not None and str(v).strip()]
                    if vals: texts.append(' | '.join(vals))
            return '\n'.join(texts)
    except Exception:
        return ''
    return ''









def table_exists(table: str) -> bool:
    try:
        validate_table_name(table)
        if database_is_persistent():
            return not query_sql("select 1 from information_schema.tables where table_schema = 'public' and table_name = :table limit 1", {'table': table}).empty
        db_all(table)
        return True
    except Exception:
        return False










def _permission_rows_for_role(role_name: str) -> pd.DataFrame:
    perms = db_all('permissions')
    rp = db_where('role_permissions', 'role_name = :role', (('role', role_name),))
    if perms.empty:
        return pd.DataFrame()
    out = perms.copy()
    enabled = set(rp[rp['enabled'] == 'Yes']['permission_id'].astype(str)) if not rp.empty else set()
    out['enabled'] = out['permission_id'].astype(str).isin(enabled)
    return out

def _effective_permission_rows(user_id: str) -> pd.DataFrame:
    users = db_all('users')
    perms = db_all('permissions')
    if users.empty or perms.empty or user_id not in set(users['user_id'].astype(str)):
        return pd.DataFrame()
    u = users[users['user_id'].astype(str) == str(user_id)].iloc[0]
    role = str(u.get('role', ''))
    base = _permission_rows_for_role(role)
    if base.empty:
        base = perms.copy()
        base['enabled'] = False
    ov = db_where('user_permission_overrides', 'user_id = :uid', (('uid', user_id),))
    if not ov.empty:
        for _, row in ov.iterrows():
            pid = str(row.get('permission_id', ''))
            if not pid:
                continue
            base.loc[base['permission_id'].astype(str) == pid, 'enabled'] = str(row.get('enabled', 'No')) == 'Yes'
    return base







def file_upload_panel(actor, linked_table='general', linked_id='general', category='Other'):
    cat = st.selectbox('File Category', FILE_CATEGORIES, index=FILE_CATEGORIES.index(category) if category in FILE_CATEGORIES else 0)
    classification = st.selectbox('Information Classification', ['Internal','Confidential','Restricted Technical','Public'], help='Controls external AI eligibility and handling of the uploaded material.')
    uploads = st.file_uploader('Upload PDF, PPT/PPTX, DOC/DOCX, TXT, images, video or Excel', type=ALLOWED_EXTENSIONS, accept_multiple_files=True)
    if st.button('Upload File(s)'):
        if not uploads:
            st.error('Select file(s).')
        else:
            count = 0
            for f in uploads:
                try:
                    upload_file(f, actor, linked_table, linked_id, cat, classification)
                    count += 1
                except Exception as e:
                    st.error(f'{f.name}: {e}')
            st.success(f'{count} file(s) uploaded.')


def phase2_health_snapshot() -> dict:
    """Return a lightweight architecture/health snapshot for admin diagnostics."""
    report = _schema_contract_report(Path(__file__).resolve().parents[1], Path(__file__).resolve().parents[1] / 'database' / 'postgres_schema.sql')
    return _application_health(database_is_persistent=database_is_persistent, storage_is_persistent=storage_is_persistent, schema_report=report) | {'performance': _performance_snapshot()}
ROLE_PERMISSION_BASELINE = {
    'GM': [('Administration','View','Organization-wide'),('Authorization','View','Organization-wide')],
    'Admin': [('Administration',a,'Organization-wide') for a in ['View','Create','Edit','Manage','Export']] + [('Users & Roles',a,'Organization-wide') for a in ['View','Create','Edit','Manage','Export']] + [('Departments',a,'Organization-wide') for a in ['View','Create','Edit','Manage','Export']] + [('Permissions',a,'Organization-wide') for a in ['View','Create','Edit','Manage','Export']] + [('System Settings',a,'Organization-wide') for a in ['View','Create','Edit','Manage','Export']] + [('Audit Trail',a,'Organization-wide') for a in ['View','Export']] + [('Backup & Recovery',a,'Organization-wide') for a in ['View','Create','Edit','Manage','Export']],
    'Trainer': [('Training',a,'Assigned') for a in ['View','Create','Edit','Assign','Review']] + [('Development Plans',a,'Assigned') for a in ['View','Create','Edit']] + [('Practical / Witness','View','Assigned'),('Competency','View','Assigned'),('Knowledge Library','View','Organization-wide')],
    'Department Manager': [('Training','View','Department'),('Practical / Witness','View','Department'),('Practical / Witness','Review','Department'),('Competency','View','Department'),('Competency','Review','Department'),('Authorization','View','Department'),('Authorization','Review','Department'),('Knowledge Library','View','Organization-wide')],
    'Surveyor': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','View','Assigned'),('Practical / Witness','Review','Assigned'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide')],
    'NSC Surveyor': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','View','Assigned'),('Practical / Witness','Review','Assigned'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide')],
    'In-Service Surveyor': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','View','Assigned'),('Practical / Witness','Review','Assigned'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide')],
    'Industrial Surveyor': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','View','Assigned'),('Practical / Witness','Review','Assigned'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide')],
    'Plan Appraiser': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','View','Assigned'),('Practical / Witness','Review','Assigned'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide')],
    'Trainee': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','Create','Own'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide'),('Development Plans','View','Own'),('Development Plans','Edit','Own')],
    'On Probation': [('Training','View','Own'),('Practical / Witness','View','Own'),('Practical / Witness','Create','Own'),('Competency','View','Own'),('Authorization','View','Own'),('Knowledge Library','View','Organization-wide'),('Development Plans','View','Own'),('Development Plans','Edit','Own')],
    'Management': [('Dashboard','View','Organization-wide'),('Authorization','View','Organization-wide'),('Authorization','Review','Organization-wide'),('Authorization','Approve','Organization-wide')],
    'QMS Auditor': [('Accreditation Readiness','View','Department'),('Knowledge Library','View','Organization-wide')],
    'QMR': [('Accreditation Readiness','View','Organization-wide'),('Knowledge Library','View','Organization-wide')],
    'Rule Development Rep': [('Interpretation Portal','View','Department'),('Knowledge Library','View','Organization-wide')],
}


# Practical/Witness professional workflow grants. Kept in the canonical baseline so
# fresh installs and upgrades share the same role/action contract.
for _role, _grants in {
    'Trainee': [('Practical / Witness','Create','Own')],
    'On Probation': [('Practical / Witness','Create','Own')],
    'Surveyor': [('Practical / Witness','View','Assigned'),('Practical / Witness','Create','Assigned'),('Practical / Witness','Review','Assigned')],
    'NSC Surveyor': [('Practical / Witness','View','Assigned'),('Practical / Witness','Create','Assigned'),('Practical / Witness','Review','Assigned')],
    'In-Service Surveyor': [('Practical / Witness','View','Assigned'),('Practical / Witness','Create','Assigned'),('Practical / Witness','Review','Assigned')],
    'Industrial Surveyor': [('Practical / Witness','View','Assigned'),('Practical / Witness','Create','Assigned'),('Practical / Witness','Review','Assigned')],
    'Plan Appraiser': [('Practical / Witness','View','Assigned'),('Practical / Witness','Create','Assigned'),('Practical / Witness','Review','Assigned')],
    'Department Manager': [('Practical / Witness','Create','Department'),('Practical / Witness','Edit','Department')],
}.items():
    ROLE_PERMISSION_BASELINE.setdefault(_role, [])
    for _grant in _grants:
        if _grant not in ROLE_PERMISSION_BASELINE[_role]:
            ROLE_PERMISSION_BASELINE[_role].append(_grant)

def _ensure_role_permission_baseline():
    perms = db_all('permissions')
    if perms.empty:
        return
    for role_name, grants in ROLE_PERMISSION_BASELINE.items():
        for module_name, action_name, scope_name in grants:
            match = perms[(perms['module_name'].astype(str) == module_name) & (perms['action'].astype(str) == action_name) & (perms['scope'].astype(str) == scope_name)]
            if match.empty:
                continue
            pid = str(match.iloc[0]['permission_id'])
            existing = db_where('role_permissions', 'role_name = :role and permission_id = :pid', (('role', role_name), ('pid', pid)))
            if existing.empty:
                db_insert('role_permissions', {'role_permission_id': uid('RPERM'), 'role_name': role_name, 'permission_id': pid, 'enabled': 'Yes', 'created_on': now(), 'updated_on': now()})
            elif str(existing.iloc[0].get('enabled', 'Yes')) != 'Yes':
                db_update('role_permissions', 'role_permission_id', existing.iloc[0]['role_permission_id'], {'enabled': 'Yes', 'updated_on': now()})



# ---------------------------------------------------------------------------
# Lazy compatibility exports
# ---------------------------------------------------------------------------
# Service modules depend on foundational helpers/constants from legacy_runtime.
# Importing those service modules back into legacy_runtime at module import time
# creates a circular import (legacy_runtime -> service -> legacy_runtime).  Keep
# the compatibility API, but resolve extracted service implementations only when
# a function is actually called, after legacy_runtime has finished initializing.

def _lazy_service_call(module_name: str, func_name: str, *args, **kwargs):
    from importlib import import_module
    fn = getattr(import_module(module_name), func_name)
    return fn(*args, **kwargs)

def create_auth_token(*a, **k): return _lazy_service_call('psb_app.services.auth_service', 'create_auth_token', *a, **k)
def resolve_auth_token(*a, **k): return _lazy_service_call('psb_app.services.auth_service', 'resolve_auth_token', *a, **k)
def clear_auth_token(*a, **k): return _lazy_service_call('psb_app.services.auth_service', 'clear_auth_token', *a, **k)
def phash(*a, **k): return _lazy_service_call('psb_app.services.auth_service', 'phash', *a, **k)
def verify_password(*a, **k): return _lazy_service_call('psb_app.services.auth_service', 'verify_password', *a, **k)
def temp_password(*a, **k): return _lazy_service_call('psb_app.services.auth_service', 'temp_password', *a, **k)

def allowed_user_ids(*a, **k): return _lazy_service_call('psb_app.services.policy_service', 'allowed_user_ids', *a, **k)
def restrict_user_frame(*a, **k): return _lazy_service_call('psb_app.services.policy_service', 'restrict_user_frame', *a, **k)
def can_action(*a, **k): return _lazy_service_call('psb_app.services.policy_service', 'can_action', *a, **k)
def access_record(*a, **k): return _lazy_service_call('psb_app.services.policy_service', 'access_record', *a, **k)

def calculate_training_progress(*a, **k): return _lazy_service_call('psb_app.services.training_service', 'calculate_training_progress', *a, **k)
def training_complete_for_user(*a, **k): return _lazy_service_call('psb_app.services.training_service', 'training_complete_for_user', *a, **k)
def get_matrix_for_scope(*a, **k): return _lazy_service_call('psb_app.services.training_service', 'get_matrix_for_scope', *a, **k)
def _training_requirement_status(*a, **k): return _lazy_service_call('psb_app.services.training_service', '_training_requirement_status', *a, **k)
def readiness(*a, **k): return _lazy_service_call('psb_app.services.training_service', 'readiness', *a, **k)
def generate_mcqs(*a, **k): return _lazy_service_call('psb_app.services.training_service', 'generate_mcqs', *a, **k)

def build_certificate(*a, **k): return _lazy_service_call('psb_app.services.certificate_service', 'build_certificate', *a, **k)

def audit(*a, **k): return _lazy_service_call('psb_app.services.governance_service', 'audit', *a, **k)
def create_notification(*a, **k): return _lazy_service_call('psb_app.services.governance_service', 'create_notification', *a, **k)
def scheduler_record(*a, **k): return _lazy_service_call('psb_app.services.governance_service', 'scheduler_record', *a, **k)
def scheduler_health_summary(*a, **k): return _lazy_service_call('psb_app.services.governance_service', 'scheduler_health_summary', *a, **k)
def kpi_definitions_frame(*a, **k): return _lazy_service_call('psb_app.services.governance_service', 'kpi_definitions_frame', *a, **k)

def _admin_only(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_admin_only', *a, **k)
def _setting_value(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_setting_value', *a, **k)
def _setting_bool(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_setting_bool', *a, **k)
def _save_setting(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_save_setting', *a, **k)
def _backup_export_tables(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_backup_export_tables', *a, **k)
def _sanitize_backup_frame(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_sanitize_backup_frame', *a, **k)
def _build_backup_payload(*a, **k): return _lazy_service_call('psb_app.services.admin_service', '_build_backup_payload', *a, **k)

def actor_get(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'actor_get', *a, **k)
def join_list(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'join_list', *a, **k)
def split_list(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'split_list', *a, **k)
def logo_data_uri(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'logo_data_uri', *a, **k)
def make_qr_data_uri(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'make_qr_data_uri', *a, **k)
def backend_status_badges(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'backend_status_badges', *a, **k)
def department_options(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'department_options', *a, **k)
def _user_label(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', '_user_label', *a, **k)
def _parse_user_label(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', '_parse_user_label', *a, **k)
def _user_label_series(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', '_user_label_series', *a, **k)
def select_person(*a, **k): return _lazy_service_call('psb_app.services.ui_helpers', 'select_person', *a, **k)

__all__ = [k for k in globals() if not k.startswith('__')]

