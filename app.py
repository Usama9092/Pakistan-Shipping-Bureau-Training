
from __future__ import annotations

from datetime import date, datetime, timedelta
import time
from pathlib import Path
from urllib.parse import quote_plus
import base64
import hashlib
import io
import json
import os
import random
import re
import secrets
import string
import uuid
import functools

import pandas as pd
import qrcode
import streamlit as st
import streamlit.components.v1 as components
from sqlalchemy import create_engine, text
from sqlalchemy.engine import Engine

try:
    from supabase import create_client
except Exception:
    create_client = None

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


APP_TITLE = "Pakistan Shipping Bureau"
APP_SUBTITLE = "World-Class Classification Society Training, Competency, Authorization and Workforce Platform"
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///psb_hrdm_world_class.db")
PUBLIC_URL = os.getenv("PUBLIC_URL", "https://training.psbureau.org")
SUPABASE_URL = os.getenv("SUPABASE_URL", "")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
SUPABASE_BUCKET = os.getenv("SUPABASE_BUCKET", "psb-hrdm-files")
LOGO_PATH = Path("assets/psb-logo.png")
LOCAL_UPLOAD_DIR = Path("local_uploads")
DEFAULT_ADMIN_PASSWORD = os.getenv("DEFAULT_ADMIN_PASSWORD", "")
MAX_LOGIN_ATTEMPTS = int(os.getenv("MAX_LOGIN_ATTEMPTS", "5"))
LOGIN_BLOCK_MINUTES = int(os.getenv("LOGIN_BLOCK_MINUTES", "5"))
SAFE_TABLE_RE = re.compile(r"^[a-zA-Z0-9_]+$")
SAFE_FILENAME_RE = re.compile(r"[^a-zA-Z0-9._-]+")
ACTIVE_SESSIONS: dict[str, dict[str, object]] = {}


# Simple in-memory rate limiter for UI actions (per-user key)
class RateLimiter:
    def __init__(self):
        # store action -> key -> list[timestamps]
        self._calls: dict[str, dict[str, list[float]]] = {}

    def _prune(self, lst: list[float], window: float) -> list[float]:
        cutoff = time.time() - window
        return [t for t in lst if t >= cutoff]

    def allowed(self, action: str, key: str, limit: int, per_seconds: int) -> bool:
        now_ts = time.time()
        self._calls.setdefault(action, {})
        self._calls[action].setdefault(key, [])
        calls = self._prune(self._calls[action][key], per_seconds)
        if len(calls) >= limit:
            return False
        calls.append(now_ts)
        self._calls[action][key] = calls
        return True


RATE_LIMITER = RateLimiter()

APP_ENV = os.getenv("APP_ENV", "production" if os.getenv("RENDER") else "local").lower()


def is_render_runtime() -> bool:
    return bool(os.getenv("RENDER") or os.getenv("RENDER_SERVICE_ID") or os.getenv("RENDER_EXTERNAL_URL"))


def database_is_persistent() -> bool:
    url = DATABASE_URL.lower().strip()
    return url.startswith(("postgresql://", "postgresql+psycopg2://", "postgres://"))


def storage_is_persistent() -> bool:
    return bool(SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY and SUPABASE_BUCKET)


def require_persistent_backend() -> None:
    """Prevent data loss on Render by blocking temporary SQLite/local storage."""
    if is_render_runtime() and not database_is_persistent():
        st.error("Persistent database is not configured. Render local SQLite storage is temporary and data will disappear after restart/redeploy.")
        st.markdown("""
        **Fix in Render → Environment Variables:**
        ```text
        DATABASE_URL=postgresql://postgres:PASSWORD@HOST:5432/postgres
        SUPABASE_URL=https://your-project.supabase.co
        SUPABASE_SERVICE_ROLE_KEY=your-service-role-key
        SUPABASE_BUCKET=psb-hrdm-files
        ```
        """)
        st.stop()


def backend_status_badges() -> str:
    db_badge = "✅ PostgreSQL/Supabase" if database_is_persistent() else "⚠️ Local SQLite"
    storage_badge = "✅ Supabase Storage" if storage_is_persistent() else ("⚠️ Local files" if not is_render_runtime() else "❌ Storage missing")
    return f"<span class='pill'>{db_badge}</span><span class='pill'>{storage_badge}</span>"


STANDARDS = [
    "IMO RO Code",
    "ISO 9001",
    "ISO/IEC 17020",
    "IACS PR7",
    "Competency-Based Qualification System",
]

ROLES = [
    "Admin",
    "Trainer",
    "Tutor/Mentor",
    "Surveyor",
    "Plan Appraiser",
    "QMS Auditor",
    "Industrial Surveyor",
    "Rule Development Rep",
    "Principal Surveyor",
    "Chief Plan Appraiser",
    "Lead Auditor",
    "Technical Manager",
    "QMR",
    "CRB Member",
    "Job Coordinator",
    "Management",
    "Trainee",
    "On Probation",
]

TRAINEE_PATHS = [
    "Trainee New Building Surveyor",
    "Trainee In-Service Surveyor",
    "Trainee Plan Appraisal Engineer",
    "Trainee QMS Auditor",
    "Trainee Industrial Surveyor",
    "Trainee Rule Development Representative",
]

JOB_TYPES = [
    "New Building Survey",
    "In-Service Survey",
    "Plan Appraisal",
    "Internal Audit",
    "External Audit",
    "Industrial Survey",
    "Rule Development",
    "Witness Survey",
]

SCOPES = [
    "Hull NB",
    "Hull IS",
    "Machinery NB",
    "Machinery IS",
    "Electrical NB",
    "Electrical IS",
    "Statutory SOLAS",
    "Statutory MARPOL",
    "Plan Approval Hull",
    "Plan Approval Machinery",
    "Plan Approval Electrical",
    "Internal Auditor",
    "External Auditor",
    "Industrial Surveyor",
    "Rule Development",
]

COMPETENCY_LEVELS = [
    "Level 0 - Trainee",
    "Level 1 - Witness Eligible",
    "Level 2 - Supervised Eligible",
    "Level 3 - Authorized",
    "Level 4 - Senior Authorized",
    "Level 5 - Principal / Lead",
]

FILE_CATEGORIES = [
    "Training Material",
    "SCORM Package",
    "Rule Document",
    "Knowledge Bulletin",
    "Survey Evidence",
    "Plan Review Evidence",
    "Witness Evidence",
    "Certificate Template",
    "Issued Certificate",
    "CAPA Evidence",
    "Other",
]

ALLOWED_EXTENSIONS = ["pdf", "ppt", "pptx", "txt", "doc", "docx", "xls", "xlsx", "png", "jpg", "jpeg", "mp4", "csv", "html"]

CORE_THEORETICAL_MODULES = [
    ("CORE-001", "PSB Induction and Code of Ethics", "All", "Core", 2),
    ("CORE-002", "IMO Recognized Organization Code Awareness", "All", "Core", 3),
    ("CORE-003", "ISO 9001 Quality Management System", "All", "QMS", 3),
    ("CORE-004", "ISO/IEC 17020 Inspection Body Requirements", "All", "QMS", 3),
    ("CORE-005", "IACS PR7 Training and Qualification Principles", "All", "Competency", 2),
    ("CORE-006", "Document Control and Record Retention", "All", "QMS", 2),
    ("CORE-007", "HSE, Risk Assessment and Site Safety", "All", "Safety", 3),
    ("CORE-008", "Survey Reporting and Deficiency Management", "Surveyor", "Survey", 3),
    ("TECH-001", "Hull Rules and Structural Survey Principles", "Hull Surveyor", "Technical", 5),
    ("TECH-002", "Machinery Rules and Machinery Survey Principles", "Machinery Surveyor", "Technical", 5),
    ("TECH-003", "Electrical Rules and Electrical Survey Principles", "Electrical Surveyor", "Technical", 5),
    ("STAT-001", "SOLAS Statutory Survey Requirements", "Statutory Surveyor", "Statutory", 5),
    ("STAT-002", "MARPOL Pollution Prevention Requirements", "Statutory Surveyor", "Statutory", 4),
    ("PLAN-001", "Plan Appraisal Rule Interpretation", "Plan Appraiser", "Plan Appraisal", 4),
    ("PLAN-002", "Plan Review Commenting and Approval Workflow", "Plan Appraiser", "Plan Appraisal", 3),
    ("AUD-001", "Internal Audit Techniques and CAPA", "Auditor", "Audit", 4),
    ("RULE-001", "Rule Development, Technical Circulars and Change Impact", "Rule Development Rep", "Rule Development", 4),
]

DEFAULT_AUTH_MATRIX = [
    ("Hull NB", "New Building Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Hull IS", "In-Service Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Machinery NB", "New Building Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Machinery IS", "In-Service Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Electrical NB", "New Building Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Electrical IS", "In-Service Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Statutory SOLAS", "Statutory Survey", 2, 1, 0, 0, 3, 3, "High", 36),
    ("Statutory MARPOL", "Statutory Survey", 2, 1, 0, 0, 3, 3, "High", 36),
    ("Plan Approval Hull", "Plan Appraisal", 0, 0, 2, 1, 3, 3, "Medium", 36),
    ("Plan Approval Machinery", "Plan Appraisal", 0, 0, 2, 1, 3, 3, "Medium", 36),
    ("Plan Approval Electrical", "Plan Appraisal", 0, 0, 2, 1, 3, 3, "Medium", 36),
    ("Internal Auditor", "Internal Audit", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("External Auditor", "External Audit", 2, 1, 0, 0, 4, 4, "High", 36),
    ("Industrial Surveyor", "Industrial Survey", 2, 1, 0, 0, 3, 3, "Medium", 36),
    ("Rule Development", "Rule Development", 1, 1, 0, 0, 4, 4, "High", 36),
]


def now() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def today() -> str:
    return date.today().strftime("%Y-%m-%d")


def uid(prefix: str) -> str:
    return f"{prefix}-{uuid.uuid4().hex[:8].upper()}"


def clean(v) -> str:
    if v is None:
        return ""
    try:
        if pd.isna(v):
            return ""
    except Exception:
        pass
    return str(v)


def validate_table_name(table: str) -> None:
    if not SAFE_TABLE_RE.fullmatch(clean(table)):
        raise ValueError("Invalid table name.")


def sanitize_path_component(value: str) -> str:
    value = clean(value).replace("\\", "/").split("/")[-1]
    value = SAFE_FILENAME_RE.sub("_", value).strip("._-")
    return value or "item"


def sanitize_filename(name: str) -> str:
    name = clean(name).replace("\\", "/").split("/")[-1]
    name = SAFE_FILENAME_RE.sub("_", name).strip("._-")
    return name or "upload"


def validate_email(email: str) -> bool:
    return bool(re.fullmatch(r"[^@ \t\r\n]+@[^@ \t\r\n]+\.[^@ \t\r\n]+", clean(email)))


def create_auth_token(user_id: str) -> str:
    token = secrets.token_urlsafe(24)
    ACTIVE_SESSIONS[token] = {"user_id": user_id, "created_on": datetime.utcnow()}
    return token


def resolve_auth_token(token: str) -> str | None:
    token = clean(token)
    session = ACTIVE_SESSIONS.get(token)
    if not session:
        return None
    created = session.get("created_on")
    if isinstance(created, datetime) and datetime.utcnow() - created > timedelta(days=1):
        ACTIVE_SESSIONS.pop(token, None)
        return None
    user_id = session.get("user_id")
    if user_id is None:
        return None
    return str(user_id)


def clear_auth_token() -> None:
    try:
        params = getattr(st, "experimental_get_query_params", lambda: {})()
        if isinstance(params, dict):
            params.pop("auth_token", None)
            getattr(st, "experimental_set_query_params", lambda **kwargs: None)(**params)
    except Exception:
        pass
    st.session_state.pop("auth_token", None)
    try:
        components.html("<script>document.cookie = 'psb_auth=; path=/; Max-Age=0;';</script>", height=0)
    except Exception:
        pass


def phash(password: str) -> str:
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def temp_password(n: int = 10) -> str:
    return "".join(secrets.choice(string.ascii_letters + string.digits + "@#$") for _ in range(n))


def days_until(date_text: str) -> int:
    if not clean(date_text):
        return 9999
    try:
        return (datetime.strptime(clean(date_text)[:10], "%Y-%m-%d").date() - date.today()).days
    except Exception:
        return 9999


def add_months(months: int) -> str:
    d = date.today()
    month = d.month - 1 + months
    year = d.year + month // 12
    month = month % 12 + 1
    day = min(d.day, [31, 29 if year % 4 == 0 else 28, 31, 30, 31, 30, 31, 31, 30, 31, 30, 31][month - 1])
    return date(year, month, day).strftime("%Y-%m-%d")


def actor_get(actor: dict, key: str, default: str = "") -> str:
    return clean(actor.get(key, default)) if isinstance(actor, dict) else default


def join_list(values: list[str]) -> str:
    return ", ".join(values)


def split_list(value: str) -> list[str]:
    return [x.strip() for x in re.split(r"[,;|]+", clean(value)) if x.strip()]


def logo_data_uri() -> str:
    if not LOGO_PATH.exists():
        return ""
    return "data:image/png;base64," + base64.b64encode(LOGO_PATH.read_bytes()).decode()


def make_qr_data_uri(value: str) -> str:
    img = qrcode.make(value)
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()


@st.cache_resource
def get_engine() -> Engine:
    url = DATABASE_URL
    if url.startswith("postgres://"):
        url = url.replace("postgres://", "postgresql+psycopg2://", 1)
    elif url.startswith("postgresql://"):
        url = url.replace("postgresql://", "postgresql+psycopg2://", 1)
    
    if url.startswith("sqlite"):
        return create_engine(url, pool_pre_ping=True, connect_args={"check_same_thread": False})
    return create_engine(
        url,
        pool_pre_ping=True,
        pool_size=int(os.getenv("DB_POOL_SIZE", "5")),
        max_overflow=int(os.getenv("DB_MAX_OVERFLOW", "10")),
        pool_recycle=1800,
        pool_timeout=30,
    )


@st.cache_resource
def get_supabase_client():
    if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY or create_client is None:
        return None
    return create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)


def exec_sql(sql: str, params: dict | None = None) -> None:
    with get_engine().begin() as conn:
        conn.execute(text(sql), params or {})


def query_sql(sql: str, params: dict | None = None) -> pd.DataFrame:
    with get_engine().begin() as conn:
        return pd.read_sql(text(sql), conn, params=params or {})


@st.cache_data(ttl=20, show_spinner=False)
def db_all(table: str) -> pd.DataFrame:
    """Cached full-table reads. Streamlit reruns the whole script on every click;
    caching prevents repeated full-table SELECTs during normal navigation.
    The cache is cleared after insert/update/delete operations below.
    """
    try:
        validate_table_name(table)
        return query_sql(f"select * from {table}")
    except Exception:
        return pd.DataFrame()

@st.cache_data(ttl=20, show_spinner=False)
def db_where(table: str, where_sql: str, params_tuple: tuple[tuple[str, object], ...] = ()) -> pd.DataFrame:
    """Cached filtered read. Use this on interactive pages instead of loading full tables."""
    try:
        validate_table_name(table)
        params = dict(params_tuple)
        return query_sql(f"select * from {table} where {where_sql}", params)
    except Exception:
        return pd.DataFrame()


def clear_db_cache() -> None:
    """Clear Streamlit data caches after write operations.
    The previous version accidentally called itself recursively, which could
    freeze the app after inserts/updates/deletes.
    """
    try:
        db_all.clear()
        db_where.clear()
    except Exception:
        pass


def first_row(df: pd.DataFrame) -> dict | None:
    if df is None or df.empty:
        return None
    return df.iloc[0].to_dict()


def convert_numpy_types(row: dict) -> dict:
    """Convert numpy types to Python native types for database compatibility."""
    converted = {}
    for key, value in row.items():
        if value is None:
            converted[key] = None
        elif hasattr(value, 'item'):  # numpy scalars have .item() method
            converted[key] = value.item()
        else:
            converted[key] = value
    return converted


def db_insert(table: str, row: dict) -> None:
    validate_table_name(table)
    row = convert_numpy_types(row)
    cols = list(row.keys())
    exec_sql(
        f"insert into {table} ({', '.join(cols)}) values ({', '.join([f':{c}' for c in cols])})",
        row,
    )
    clear_db_cache()


def db_update(table: str, id_col: str, id_val: str, row: dict) -> None:
    if not row:
        return
    validate_table_name(table)
    patch = dict(row)
    patch[id_col] = id_val
    patch = convert_numpy_types(patch)
    sets = ", ".join([f"{k}=:{k}" for k in row.keys()])
    exec_sql(f"update {table} set {sets} where {id_col}=:{id_col}", patch)
    clear_db_cache()


def db_delete(table: str, id_col: str, id_val: str) -> None:
    validate_table_name(table)
    exec_sql(f"delete from {table} where {id_col} = :id", {"id": id_val})
    clear_db_cache()


@st.cache_resource(show_spinner=False)
def init_db() -> None:
    stmts = [
        """create table if not exists users (
            user_id text primary key, name text, role text, trainee_path text, department text, assigned_duty text,
            email text unique, login_id text unique, password_hash text, temp_password text, status text,
            availability text, current_location text, mentor_id text, mentor_name text, competency_level text,
            created_on text, last_login text
        )""",
        """create table if not exists training_modules (
            module_id text primary key, title text, module_group text, target_path text, mandatory text,
            refresher_required text, cpd_hours real, validity_months integer, added_by text, created_on text
        )""",
        """create table if not exists trainings (
            training_id text primary key, module_id text, title text, category text, standards text, target_roles text,
            target_paths text, trainer_id text, trainer_name text, slides_link text, video_link text, reference_link text,
            scorm_package_link text, lms_course_id text, schedule_date text, schedule_time text, meeting_link text,
            recording_link text, passing_marks integer, validity_months integer, max_attempts integer, retest_wait_days integer,
            status text, created_on text, updated_on text
        )""",
        """create table if not exists files (
            file_id text primary key, owner_user_id text, owner_name text, linked_table text, linked_id text,
            category text, file_name text, file_ext text, mime_type text, storage_provider text,
            storage_path text, public_url text, extracted_text text, ocr_status text, review_status text,
            created_on text, updated_on text
        )""",
        """create table if not exists training_records (
            record_id text primary key, user_id text, name text, role text, trainee_path text, training_id text,
            training_title text, status text, slides_opened text, video_opened text, live_attendance text,
            recording_opened text, lms_completed text, test_status text, score real, passing_marks integer,
            certificate_status text, certificate_link text, due_date text, completed_on text, progress integer,
            remarks text, updated_on text
        )""",
        """create table if not exists question_bank (
            question_id text primary key, training_id text, question text, option_a text, option_b text,
            option_c text, option_d text, correct_answer text, marks integer, generated_on text
        )""",
        """create table if not exists assessment_history (
            assessment_id text primary key, user_id text, name text, training_id text, training_title text,
            attempt_no integer, score real, result text, attempted_on text, next_retest_allowed text, remarks text
        )""",
        """create table if not exists competency_matrix (
            competency_id text primary key, user_id text, name text, role text, trainee_path text, area text,
            competency_level text, scope text, job_type text, required_training_ids text, required_witness_count integer,
            required_supervised_count integer, required_joint_plan_count integer, required_independent_plan_count integer,
            required_level_for_auth text, status text, expiry_date text, evidence text, created_on text, updated_on text
        )""",
        """create table if not exists authorization_matrix (
            matrix_id text primary key, scope text, job_type text, required_witness_count integer,
            required_supervised_count integer, required_joint_plan_count integer, required_independent_plan_count integer,
            required_level_for_auth text, minimum_job_level text, risk_category text, validity_months integer, active text
        )""",
        """create table if not exists development_plans (
            plan_id text primary key, user_id text, name text, trainee_path text, mentor_id text, mentor_name text,
            competency_scope text, month_no integer, activity text, target_date text, status text, mentor_comments text,
            created_on text, updated_on text
        )""",
        """create table if not exists field_exposure_matrix (
            exposure_id text primary key, user_id text, name text, trainee_path text, scope text, activity_type text,
            required_count integer, completed_count integer, status text, updated_on text
        )""",
        """create table if not exists witness_surveys (
            witness_id text primary key, user_id text, name text, trainee_path text, tutor_id text, tutor_name text,
            vessel_or_project text, job_type text, scope text, witness_date text, location text, technical_knowledge integer,
            rule_application integer, safety_awareness integer, communication integer, report_quality integer,
            professional_conduct integer, outcome text, comments text, status text, created_on text, updated_on text
        )""",
        """create table if not exists supervised_activities (
            supervised_id text primary key, user_id text, name text, trainee_path text, tutor_id text, tutor_name text,
            activity_kind text, vessel_or_project text, job_type text, scope text, activity_date text, location text,
            preparation integer, execution_quality integer, findings_quality integer, reporting_quality integer,
            rule_compliance integer, outcome text, comments text, status text, created_on text, updated_on text
        )""",
        """create table if not exists authorization_requests (
            authorization_id text primary key, user_id text, name text, trainee_path text, job_type text, scope text,
            competency_id text, status text, tutor_remarks text, tutor_signature text, tutor_signed_on text,
            principal_remarks text, principal_signature text, principal_signed_on text, technical_remarks text,
            technical_signature text, technical_signed_on text, qms_remarks text, qms_signature text, qms_signed_on text,
            crb_decision text, crb_remarks text, management_remarks text, management_signature text,
            management_signed_on text, expiry_date text, certificate_id text, certificate_html text,
            certificate_storage_link text, qr_data_uri text, created_on text, updated_on text
        )""",
        """create table if not exists authorization_certificates (
            certificate_id text primary key, authorization_id text, user_id text, name text, scope text, job_type text,
            issue_date text, expiry_date text, certificate_html text, qr_data_uri text, storage_link text,
            verification_url text, status text, created_on text
        )""",
        """create table if not exists crb_reviews (
            crb_id text primary key, authorization_id text, user_id text, name text, scope text, review_date text,
            tutor_decision text, technical_decision text, qmr_decision text, management_decision text,
            final_decision text, remarks text, signed_by text, created_on text
        )""",
        """create table if not exists annual_reviews (
            review_id text primary key, user_id text, name text, scope text, review_year integer,
            training_status text, kpi_status text, complaint_status text, capa_status text, decision text,
            reviewer text, review_date text, remarks text
        )""",
        """create table if not exists revalidation_requests (
            revalidation_id text primary key, authorization_id text, user_id text, name text, scope text,
            refresher_training_status text, annual_review_status text, kpi_review_status text, tutor_confirmation text,
            crb_status text, final_status text, due_date text, created_on text, updated_on text
        )""",
        """create table if not exists job_requests (
            job_id text primary key, job_title text, job_type text, required_scope text, vessel_name text,
            imo_number text, location text, planned_date text, priority text, risk_level text, minimum_level text,
            status text, created_by text, assigned_user_id text, assigned_user_name text, assignment_reason text,
            created_on text, updated_on text
        )""",
        """create table if not exists kpi_records (
            kpi_id text primary key, user_id text, name text, period text, surveys_done integer,
            plans_reviewed integer, audits_done integer, reports_overdue integer, ncr_count integer,
            client_feedback real, training_compliance real, utilization_percent real, kpi_score real,
            created_on text, remarks text
        )""",
        """create table if not exists cpd_records (
            cpd_id text primary key, user_id text, name text, title text, category text, hours real,
            provider text, completion_date text, evidence_file_id text, status text, created_on text
        )""",
        """create table if not exists knowledge_library (
            knowledge_id text primary key, title text, category text, standard text, revision text, issue_date text,
            file_id text, mandatory_ack text, uploaded_by text, created_on text
        )""",
        """create table if not exists knowledge_acknowledgements (
            ack_id text primary key, knowledge_id text, user_id text, name text, acknowledged_on text, status text
        )""",
        """create table if not exists rule_library (
            rule_id text primary key, title text, standard text, revision text, category text, link text,
            mandatory text, current_version_id text, created_on text, updated_on text
        )""",
        """create table if not exists document_versions (
            version_id text primary key, rule_id text, version_no text, revision_date text, change_summary text,
            file_link text, uploaded_by text, approved_by text, status text, created_on text
        )""",
        """create table if not exists capa_register (
            capa_id text primary key, source text, finding text, severity text, owner_id text, owner_name text,
            due_date text, status text, corrective_action text, created_on text, updated_on text
        )""",
        """create table if not exists notifications (
            notification_id text primary key, user_id text, name text, email text, subject text, message text,
            type text, status text, created_on text, sent_on text
        )""",
        """create table if not exists audit_trail (
            audit_id text primary key, date_time text, actor_id text, actor_name text, actor_role text,
            action text, details text, result text
        )""",
        """create table if not exists technical_authorities (
            authority_id text primary key, user_id text, name text, discipline text, authority_level text,
            approval_limit text, active text, appointed_by text, appointed_on text, remarks text
        )""",
        """create table if not exists survey_report_reviews (
            review_id text primary key, user_id text, name text, survey_scope text, vessel_name text,
            report_file_id text, reviewer_id text, reviewer_name text, technical_quality integer,
            deficiency_identification integer, rule_interpretation integer, report_writing integer,
            decision_quality integer, overall_score real, decision text, comments text, created_on text
        )""",
        """create table if not exists plan_review_quality (
            planqa_id text primary key, user_id text, name text, plan_scope text, project_name text,
            plan_file_id text, reviewer_id text, reviewer_name text, comments_quality integer,
            missed_findings integer, turnaround_days integer, accuracy_score integer, overall_score real,
            result text, comments text, created_on text
        )""",
        """create table if not exists competency_ncrs (
            ncr_id text primary key, user_id text, name text, source text, scope text, ncr_type text,
            description text, severity text, impact_on_authorization text, status text, corrective_action text,
            raised_by text, raised_on text, closed_on text
        )""",
        """create table if not exists authorization_restrictions (
            restriction_id text primary key, authorization_id text, user_id text, name text, scope text,
            restriction_type text, restriction_detail text, effective_date text, expiry_date text, status text,
            imposed_by text, created_on text
        )""",
        """create table if not exists client_feedback (
            feedback_id text primary key, user_id text, name text, client_name text, project_or_vessel text,
            job_id text, rating integer, feedback_type text, comments text, impact_on_kpi text, received_on text
        )""",
        """create table if not exists succession_plans (
            succession_id text primary key, user_id text, name text, current_role_name text, target_role text,
            readiness_level text, successor_for text, development_actions text, expected_ready_date text,
            sponsor text, status text, created_on text
        )""",
        """create table if not exists workforce_forecasts (
            forecast_id text primary key, forecast_period text, discipline text, required_headcount integer,
            available_headcount integer, expiring_authorizations integer, leave_or_unavailable integer,
            gap integer, risk_status text, mitigation_plan text, created_on text
        )""",
        """create table if not exists accreditation_evidence (
            evidence_id text primary key, standard text, clause text, requirement text, linked_table text,
            linked_id text, evidence_summary text, status text, owner text, last_reviewed text
        )""",
        """create table if not exists technical_interpretations (
            interpretation_id text primary key, title text, discipline text, related_rule text, question text,
            interpretation text, approved_by text, approval_status text, revision text, issue_date text,
            created_on text
        )""",

    ]
    for s in stmts:
        exec_sql(s)
    ensure_indexes()
    if db_all("users").empty:
        seed_demo()


def ensure_indexes() -> None:
    """Create common PostgreSQL/Supabase indexes used by dashboards and trainee pages."""
    indexes = [
        "create index if not exists users_login_id_idx on users(login_id)",
        "create index if not exists users_email_idx on users(email)",
        "create index if not exists trainings_trainer_id_idx on trainings(trainer_id)",
        "create index if not exists trainings_status_idx on trainings(status)",
        "create index if not exists training_records_user_id_idx on training_records(user_id)",
        "create index if not exists training_records_training_id_idx on training_records(training_id)",
        "create index if not exists training_records_user_training_idx on training_records(user_id, training_id)",
        "create index if not exists files_owner_user_id_idx on files(owner_user_id)",
        "create index if not exists files_linked_idx on files(linked_table, linked_id)",
        "create index if not exists notifications_user_id_idx on notifications(user_id)",
        "create index if not exists question_bank_training_id_idx on question_bank(training_id)",
        "create index if not exists assessment_history_user_training_idx on assessment_history(user_id, training_id)",
        "create index if not exists competency_matrix_user_id_idx on competency_matrix(user_id)",
        "create index if not exists authorization_requests_user_id_idx on authorization_requests(user_id)",
        "create index if not exists job_requests_assigned_user_id_idx on job_requests(assigned_user_id)",
        "create index if not exists kpi_records_user_id_idx on kpi_records(user_id)",
        "create index if not exists cpd_records_user_id_idx on cpd_records(user_id)",
    ]
    for idx in indexes:
        try:
            exec_sql(idx)
        except Exception:
            pass


def audit(action: str, details: str | None = "", result: str = "Success", actor: dict | None = None) -> None:
    actor_data = actor or st.session_state.get("user", {})
    details = clean(details)
    db_insert("audit_trail", {
        "audit_id": uid("AUD"),
        "date_time": now(),
        "actor_id": actor_get(actor_data, "user_id"),
        "actor_name": actor_get(actor_data, "name", "System"),
        "actor_role": actor_get(actor_data, "role", "System"),
        "action": action,
        "details": details,
        "result": result,
    })


def seed_demo() -> None:
    demo_users = [
        ("USR-ADMIN", "PSB Admin", "Admin", "", "Support/Admin", "System Control", "admin@psbureau.org", "admin", "Admin@1234", "", ""),
        ("USR-MGMT", "Management User", "Management", "", "Management", "Oversight", "management@psbureau.org", "management", "Mgmt@1234", "", ""),
        ("USR-TRAINER", "Training Officer", "Trainer", "", "Training", "Training Delivery", "trainer@psbureau.org", "trainer", "Trainer@1234", "", ""),
        ("USR-TUTOR", "Senior Surveyor Tutor", "Tutor/Mentor", "", "Survey", "Mentor and Witness Evaluation", "tutor@psbureau.org", "tutor", "Tutor@1234", "", ""),
        ("USR-TECH", "Technical Manager", "Technical Manager", "", "Technical", "Technical Authority", "technical@psbureau.org", "technical", "Tech@1234", "", ""),
        ("USR-PRINCIPAL", "Principal Surveyor", "Principal Surveyor", "", "Survey", "Principal Authority", "principal@psbureau.org", "principal", "Principal@1234", "", ""),
        ("USR-QMR", "QMS Representative", "QMR", "", "QMS", "QMS Review", "qmr@psbureau.org", "qmr", "QMR@1234", "", ""),
        ("USR-COORD", "Job Coordinator", "Job Coordinator", "", "Operations", "Job Allocation", "coordinator@psbureau.org", "coordinator", "Coord@1234", "", ""),
        ("USR-SURVEYOR", "Sample Trainee Surveyor", "Trainee", "Trainee New Building Surveyor", "Survey", "Electrical NB Path", "surveyor@psbureau.org", "surveyor", "Surveyor@1234", "USR-TUTOR", "Senior Surveyor Tutor"),
        ("USR-APPRAISER", "Sample Trainee Plan Appraiser", "Trainee", "Trainee Plan Appraisal Engineer", "Plan Appraisal", "Electrical Plan Approval Path", "appraiser@psbureau.org", "appraiser", "Appraiser@1234", "USR-TUTOR", "Senior Surveyor Tutor"),
    ]
    for u in demo_users:
        password = DEFAULT_ADMIN_PASSWORD or temp_password(12)
        db_insert("users", {
            "user_id": u[0], "name": u[1], "role": u[2], "trainee_path": u[3], "department": u[4],
            "assigned_duty": u[5], "email": u[6], "login_id": u[7], "password_hash": phash(password),
            "temp_password": "", "status": "Active", "availability": "Available", "current_location": "Karachi",
            "mentor_id": u[9], "mentor_name": u[10], "competency_level": "Level 0 - Trainee",
            "created_on": today(), "last_login": "",
        })
    for module in CORE_THEORETICAL_MODULES:
        db_insert("training_modules", {
            "module_id": module[0], "title": module[1], "module_group": module[3], "target_path": module[2],
            "mandatory": "Yes", "refresher_required": "Yes", "cpd_hours": module[4], "validity_months": 36,
            "added_by": "System", "created_on": today(),
        })
    for row in DEFAULT_AUTH_MATRIX:
        db_insert("authorization_matrix", {
            "matrix_id": uid("MATRIX"), "scope": row[0], "job_type": row[1],
            "required_witness_count": row[2], "required_supervised_count": row[3],
            "required_joint_plan_count": row[4], "required_independent_plan_count": row[5],
            "required_level_for_auth": f"Level {row[6]} - Authorized" if row[6] == 3 else f"Level {row[6]} - Senior Authorized",
            "minimum_job_level": f"Level {row[7]} - Authorized" if row[7] == 3 else f"Level {row[7]} - Senior Authorized",
            "risk_category": row[8], "validity_months": row[9], "active": "Yes",
        })
    for rule in [
        ("RULE-IMO-RO", "IMO Recognized Organization Code", "IMO RO Code", "Current", "Statutory", "https://www.imo.org"),
        ("RULE-ISO9001", "Quality Management System Requirements", "ISO 9001", "2015", "QMS", "https://www.iso.org"),
        ("RULE-ISO17020", "Inspection Body Competence Requirements", "ISO/IEC 17020", "2012", "Inspection", "https://www.iso.org"),
        ("RULE-IACS-PR7", "IACS Training and Qualification Principles", "IACS PR7", "Current", "Competency", "https://iacs.org.uk"),
    ]:
        db_insert("rule_library", {
            "rule_id": rule[0], "title": rule[1], "standard": rule[2], "revision": rule[3],
            "category": rule[4], "link": rule[5], "mandatory": "Yes", "current_version_id": "",
            "created_on": today(), "updated_on": today(),
        })
    audit("Database Seeded", "World-class PSB HRDM data seeded", actor={"name": "System", "role": "System"})


def upload_file(uploaded_file, actor: dict, linked_table: str, linked_id: str, category: str) -> dict:
    # basic per-user rate limiting for uploads
    try:
        user_key = actor_get(actor, "user_id") or "anon"
        if not RATE_LIMITER.allowed("upload", user_key, limit=15, per_seconds=60):
            raise RuntimeError("Rate limit exceeded for uploads. Try again later.")
    except Exception:
        # on any limiter error, allow to continue (do not block app startup) or re-raise
        pass
    file_id = uid("FILE")
    ext = uploaded_file.name.split(".")[-1].lower() if "." in uploaded_file.name else ""
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"File type .{ext} is not allowed.")
    data = uploaded_file.getvalue()
    linked_table = sanitize_path_component(linked_table)
    linked_id = sanitize_path_component(linked_id)
    filename = sanitize_filename(uploaded_file.name)
    storage_path = f"{sanitize_path_component(category)}/{linked_table}/{linked_id}/{file_id}_{filename}"
    provider = "local"
    public_url = ""
    client = get_supabase_client()
    if client is not None:
        try:
            try:
                client.storage.create_bucket(SUPABASE_BUCKET, options={"public": True})
            except Exception:
                pass
            client.storage.from_(SUPABASE_BUCKET).upload(
                storage_path, data,
                {"content-type": uploaded_file.type or "application/octet-stream", "upsert": "true"}
            )
            public_url = client.storage.from_(SUPABASE_BUCKET).get_public_url(storage_path)
            provider = "supabase"
        except Exception as e:
            if is_render_runtime():
                raise RuntimeError(f"Supabase Storage upload failed on Render. Configure SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY and SUPABASE_BUCKET. Details: {e}")
            provider = "local"
    if provider == "local":
        if is_render_runtime():
            raise RuntimeError("Local file storage is disabled on Render because it is temporary. Configure Supabase Storage.")
        local_path = LOCAL_UPLOAD_DIR / storage_path
        local_path.parent.mkdir(parents=True, exist_ok=True)
        local_path.write_bytes(data)
        public_url = str(local_path)
    extracted = extract_text(uploaded_file.name, data)
    row = {
        "file_id": file_id, "owner_user_id": actor_get(actor, "user_id"), "owner_name": actor_get(actor, "name"),
        "linked_table": linked_table, "linked_id": linked_id, "category": category, "file_name": uploaded_file.name,
        "file_ext": ext, "mime_type": uploaded_file.type or "", "storage_provider": provider,
        "storage_path": storage_path, "public_url": public_url, "extracted_text": extracted[:10000],
        "ocr_status": "Extracted" if extracted else "Pending/Not Supported", "review_status": "Pending Review",
        "created_on": now(), "updated_on": now(),
    }
    db_insert("files", row)
    audit("File Uploaded", f"{uploaded_file.name} linked to {linked_table}:{linked_id}", actor=actor)
    return row


def extract_text(name: str, data: bytes) -> str:
    lower = name.lower()
    try:
        if lower.endswith((".txt", ".csv")):
            return data.decode("utf-8", errors="ignore")
        if lower.endswith(".pdf") and PdfReader:
            reader = PdfReader(io.BytesIO(data))
            return "\n".join([p.extract_text() or "" for p in reader.pages])
        if lower.endswith(".docx") and docx:
            doc = docx.Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)
        if lower.endswith(".pptx") and Presentation:
            prs = Presentation(io.BytesIO(data))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)  # type: ignore
            return "\n".join(texts)
    except Exception:
        return ""
    return ""


def create_notification(user_id: str, subject: str, message: str, ntype: str) -> None:
    u = db_where("users", "user_id = :user_id", (("user_id", user_id),))
    if u.empty:
        return
    row = u.iloc[0]
    db_insert("notifications", {
        "notification_id": uid("NOT"), "user_id": row["user_id"], "name": row["name"], "email": row["email"],
        "subject": subject, "message": message, "type": ntype, "status": "Generated", "created_on": now(), "sent_on": "",
    })


def calculate_training_progress(r: pd.Series) -> tuple[int, str, str]:
    checks = [
        r.get("slides_opened") == "Yes",
        r.get("video_opened") == "Yes" or r.get("recording_opened") == "Yes",
        r.get("live_attendance") in ["Present", "Recording Viewed"],
        r.get("lms_completed") == "Yes",
        r.get("test_status") == "Passed",
        r.get("certificate_status") == "Issued",
    ]
    progress = int(sum(checks) / len(checks) * 100)
    status = "Completed" if progress == 100 else "Pending"
    completed_on = today() if progress == 100 and not clean(r.get("completed_on")) else clean(r.get("completed_on"))
    return progress, status, completed_on


def update_training_progress(record_id: str | None = None) -> None:
    """Update one record where possible. Full-table updates made each click very slow."""
    if record_id:
        records = db_where("training_records", "record_id = :record_id", (("record_id", record_id),))
    else:
        records = db_all("training_records")
    for _, r in records.iterrows():
        progress, status, completed_on = calculate_training_progress(r)
        patch = {"progress": progress, "status": status, "completed_on": completed_on, "updated_on": now()}
        db_update("training_records", "record_id", r["record_id"], patch)


def training_complete_for_user(user_id: str) -> bool:
    assigned = db_where("training_records", "user_id = :user_id", (("user_id", user_id),))
    return not assigned.empty and len(assigned[assigned["test_status"] != "Passed"]) == 0


def get_matrix_for_scope(scope: str) -> pd.Series | None:
    matrix = db_all("authorization_matrix")
    m = matrix[(matrix["scope"] == scope) & (matrix["active"] == "Yes")] if not matrix.empty else pd.DataFrame()
    if m.empty:
        return None
    return m.iloc[0]


def readiness(user_id: str, scope: str) -> tuple[bool, list[str]]:
    matrix = get_matrix_for_scope(scope)
    if matrix is None:
        return False, ["No authorization matrix defined for this scope."]
    gaps = []
    records = db_all("training_records")
    if records.empty or records[(records["user_id"] == user_id) & (records["test_status"] == "Passed")].empty:
        gaps.append("No passed theoretical training / assessment found.")
    witness = db_all("witness_surveys")
    witness_count = len(witness[(witness["user_id"] == user_id) & (witness["scope"] == scope) & (witness["outcome"] == "Pass")]) if not witness.empty else 0
    if witness_count < int(matrix["required_witness_count"]):
        gaps.append(f"Witness surveys incomplete: {witness_count}/{matrix['required_witness_count']}.")
    sup = db_all("supervised_activities")
    sup_count = len(sup[(sup["user_id"] == user_id) & (sup["scope"] == scope) & (sup["activity_kind"].isin(["Supervised Survey", "Independent Audit", "Supervised Rule Exercise"])) & (sup["outcome"] == "Pass")]) if not sup.empty else 0
    if sup_count < int(matrix["required_supervised_count"]):
        gaps.append(f"Supervised activities incomplete: {sup_count}/{matrix['required_supervised_count']}.")
    joint_count = len(sup[(sup["user_id"] == user_id) & (sup["scope"] == scope) & (sup["activity_kind"] == "Joint Plan Review") & (sup["outcome"] == "Pass")]) if not sup.empty else 0
    if joint_count < int(matrix["required_joint_plan_count"]):
        gaps.append(f"Joint plan reviews incomplete: {joint_count}/{matrix['required_joint_plan_count']}.")
    indep_count = len(sup[(sup["user_id"] == user_id) & (sup["scope"] == scope) & (sup["activity_kind"] == "Independent Plan Review") & (sup["outcome"] == "Pass")]) if not sup.empty else 0
    if indep_count < int(matrix["required_independent_plan_count"]):
        gaps.append(f"Independent plan reviews incomplete: {indep_count}/{matrix['required_independent_plan_count']}.")
    plans = db_all("development_plans")
    open_plan = len(plans[(plans["user_id"] == user_id) & (plans["status"] != "Completed")]) if not plans.empty else 0
    if open_plan > 0:
        gaps.append(f"Development plan has {open_plan} open item(s).")
    return len(gaps) == 0, gaps


def generate_mcqs(training_id: str, text_value: str, count: int) -> pd.DataFrame:
    stop = {"training","system","should","shall","which","there","their","about","through","during","after","before","within","using","based","these","those","where","under","requirements","procedure","document","classification","society","survey","surveyor","appraisal","management","development"}
    keys = []
    for w in re.findall(r"\b[A-Za-z][A-Za-z\-]{4,}\b", text_value):
        x = w.lower()
        t = x.title()
        if x not in stop and t not in keys:
            keys.append(t)
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text_value.replace("\n", " ")) if 45 <= len(s.strip()) <= 260]
    rows = []
    if len(keys) < 4:
        return pd.DataFrame()
    random.shuffle(sentences)
    for s in sentences:
        if len(rows) >= count:
            break
        ans = next((k for k in keys if re.search(rf"\b{re.escape(k)}\b", s, re.I)), None)
        if not ans:
            continue
        distractors = [k for k in keys if k.lower() != ans.lower()]
        if len(distractors) < 3:
            continue
        opts = random.sample(distractors, 3) + [ans]
        random.shuffle(opts)
        rows.append({
            "question_id": uid("Q"), "training_id": training_id,
            "question": re.sub(rf"\b{re.escape(ans)}\b", "__________", s, flags=re.I, count=1),
            "option_a": opts[0], "option_b": opts[1], "option_c": opts[2], "option_d": opts[3],
            "correct_answer": ans, "marks": 1, "generated_on": now(),
        })
    return pd.DataFrame(rows)


def build_certificate(auth: pd.Series) -> tuple[str, str, str]:
    cert_id = clean(auth.get("certificate_id")) or uid("CERT")
    verification_url = f"{PUBLIC_URL}/verify/{cert_id}"
    qr = make_qr_data_uri(verification_url)
    html = f"""
<!doctype html>
<html><head><meta charset='utf-8'><title>PSB Authorization Certificate</title>
<style>
body{{font-family:Arial,sans-serif;padding:40px;color:#0f172a;background:transparent}}
.cert{{border:5px solid #071225;padding:35px;border-radius:18px;background:#fff}}
h1{{color:#071225;text-align:center;margin-bottom:0}} h2{{text-align:center;color:#0b3b76;margin-top:6px}}
.row{{margin:12px 0;font-size:16px}} .sig{{display:grid;grid-template-columns:repeat(4,1fr);gap:16px;margin-top:45px}}
.box{{border-top:1px solid #0f172a;padding-top:8px;font-size:12px}} .qr{{text-align:center;margin-top:25px}}
@media (prefers-color-scheme: dark) {{
    body{{color:#e6eef8}}
    .cert{{background:#071122;border-color:rgba(255,255,255,0.06)}}
    h1{{color:#e6eef8}} h2{{color:#bfe0ff}}
    .box{{border-top:1px solid rgba(255,255,255,0.06)}}
}}
</style></head><body><div class='cert'>
<h1>Pakistan Shipping Bureau</h1><h2>Authorization Certificate / Letter</h2>
<div class='row'><b>Certificate ID:</b> {cert_id}</div>
<div class='row'><b>Authorization ID:</b> {auth['authorization_id']}</div>
<div class='row'><b>Name:</b> {auth['name']}</div>
<div class='row'><b>Role/Path:</b> {auth['trainee_path']}</div>
<div class='row'><b>Job Type:</b> {auth['job_type']}</div>
<div class='row'><b>Authorized Scope:</b> {auth['scope']}</div>
<div class='row'><b>Status:</b> {auth['status']}</div>
<div class='row'><b>Valid Until:</b> {auth['expiry_date']}</div>
<div class='row'><b>Standards Basis:</b> {", ".join(STANDARDS)}</div>
<div class='sig'>
<div class='box'><b>Tutor</b><br>{auth.get('tutor_signature','')}</div>
<div class='box'><b>Principal/Chief</b><br>{auth.get('principal_signature','')}</div>
<div class='box'><b>QMS</b><br>{auth.get('qms_signature','')}</div>
<div class='box'><b>Management</b><br>{auth.get('management_signature','')}</div>
</div>
<div class='qr'><img src='{qr}' width='125'><br><small>Verify: {verification_url}</small></div>
</div></body></html>
"""
    return cert_id, html, qr


def apply_style() -> None:
    st.markdown("""
    <style>
    :root{--psb-navy:#071225;--psb-blue:#0b3b76;--psb-sky:#124f9e;--psb-card:#ffffff;--psb-line:#dbe3ef;--psb-text:#0f172a;--psb-muted:#64748b}
    @media (prefers-color-scheme: dark) {
        :root{--psb-card:#0b1220;--psb-line:rgba(255,255,255,0.06);--psb-text:#e6eef8;--psb-muted:#9aa7b8}
        .stApp{background:radial-gradient(circle at top left,#041022 0,#07121a 34%,#052036 100%);color:var(--psb-text)}
        section[data-testid="stSidebar"]{background:linear-gradient(180deg,var(--psb-navy) 0%,var(--psb-blue) 72%,#08244b 100%);border-right:1px solid rgba(255,255,255,.03)}
        .psb-card, .step, div[data-testid="stMetric"]{background:var(--psb-card);border:1px solid var(--psb-line)}
        /* Panels, cards and form controls */
        .login-panel{background:transparent}
        .login-frame{background:linear-gradient(180deg,rgba(6,10,18,0.6),rgba(4,10,20,0.65));border-color:rgba(255,255,255,0.03);box-shadow:0 30px 80px rgba(0,0,0,0.6)}
        .login-card{background:var(--psb-card);border:1px solid var(--psb-line);box-shadow:0 14px 40px rgba(2,6,12,0.6)}
        .login-card h2{color:var(--psb-text)}
        .login-card .muted{color:var(--psb-muted)}
        .login-card label{color:var(--psb-text)!important}
        input, textarea, select, .stTextInput>div>div input, .stTextArea>div>div textarea{background:var(--psb-card)!important;color:var(--psb-text)!important;border:1px solid var(--psb-line)!important}
        input::placeholder, textarea::placeholder{color:var(--psb-muted)!important}
        .stButton>button, .stDownloadButton>button{background:linear-gradient(135deg,var(--psb-blue),var(--psb-navy));border-color:var(--psb-navy);color:white}
        .stExpander, div[data-testid="stExpander"]{background:transparent;border:1px solid var(--psb-line)}
        div[data-testid="stDataFrame"], .stTabs [data-baseweb="tab"], .stTabs [data-baseweb="tab-list"]{background:var(--psb-card)!important;border:1px solid var(--psb-line)}
        h1,h2,h3{color:var(--psb-text)}
    }
    .stApp{background:radial-gradient(circle at top left,#eaf2ff 0,#f8fafc 34%,#eef3f8 100%);color:var(--psb-text)}
    /* Inputs, controls and buttons use card/bg vars so they adapt to theme */
    input, textarea, select, button, .stButton > button, .stTextInput>div>div, .stTextArea>div>div{background:var(--psb-card)!important;color:var(--psb-text)!important;border:1px solid var(--psb-line)!important}
    .stButton>button{box-shadow:none;border-radius:10px;padding:8px 12px}
    a, a:hover{color:var(--psb-sky)}
    .block-container{padding-top:1rem;padding-bottom:2.5rem;max-width:1480px}
    #MainMenu, footer, header[data-testid="stHeader"]{visibility:hidden}
    button[title="Toggle sidebar"], button[aria-label="Toggle sidebar"], button[aria-label="Collapse sidebar"], button[aria-label="Expand sidebar"], div[role="button"][aria-label*="sidebar"]{display:none!important}
    section[data-testid="stSidebar"]{background:linear-gradient(180deg,var(--psb-navy) 0%,var(--psb-blue) 72%,#08244b 100%);border-right:1px solid rgba(255,255,255,.10);visibility:visible!important}
    section[data-testid="stSidebar"] *{color:#f8fafc}
    section[data-testid="stSidebar"] [data-testid="stRadio"] label{font-weight:800;letter-spacing:.02em}
    section[data-testid="stSidebar"] div[role="radiogroup"] label{border-radius:12px;padding:.35rem .55rem;margin:.12rem 0}
    section[data-testid="stSidebar"] div[role="radiogroup"] label:hover{background:rgba(255,255,255,.11)}
    div[data-testid="stMetric"]{background:var(--psb-card);border:1px solid var(--psb-line);border-radius:20px;padding:16px;box-shadow:0 14px 35px rgba(15,23,42,.08)}
    div[data-testid="stMetric"] label{color:var(--psb-muted)!important;font-weight:700}
    .psb-hero{background:linear-gradient(135deg,var(--psb-navy),var(--psb-blue) 62%,var(--psb-sky));color:white;padding:1.55rem 1.85rem;border-radius:30px;margin-bottom:1.3rem;box-shadow:0 26px 75px rgba(15,23,42,.25);display:flex;gap:20px;align-items:center;border:1px solid rgba(255,255,255,.17)}
    .psb-hero img{width:96px;height:96px;border-radius:50%;object-fit:contain;background:white;padding:6px;box-shadow:0 14px 34px rgba(0,0,0,.25)}
    .psb-hero h1{margin:0;font-size:2.18rem;letter-spacing:-.035em;font-weight:900}
    .psb-hero p{color:#dbeafe;margin:.42rem 0 .25rem;font-size:1.03rem}
    .pill{display:inline-flex;padding:6px 12px;border-radius:999px;background:#e8eef7;color:#0f172a;font-size:12px;font-weight:800;margin:4px 5px 4px 0;border:1px solid #d7e0ec;white-space:nowrap}
    .psb-hero .pill{background:rgba(255,255,255,.14);color:white;border:1px solid rgba(255,255,255,.24)}
    .step{border-left:5px solid var(--psb-blue);background:white;border-radius:18px;padding:.9rem 1rem;margin:.48rem 0;box-shadow:0 12px 32px rgba(15,23,42,.07)}
    .psb-card{background:white;border:1px solid var(--psb-line);border-radius:22px;padding:1rem 1.1rem;margin:.65rem 0;box-shadow:0 12px 32px rgba(15,23,42,.07)}
    .psb-section-title{font-size:1.02rem;font-weight:900;color:var(--psb-blue);margin:.25rem 0 .65rem}
    .login-shell{min-height:calc(100vh - 3.5rem);display:flex;align-items:center;justify-content:center;padding:1.5rem 0 2.8rem}
    .login-frame{width:min(1180px,96vw);display:grid;grid-template-columns:1.08fr .92fr;gap:0;background:rgba(255,255,255,.84);border:1px solid rgba(219,227,239,.95);border-radius:36px;overflow:hidden;box-shadow:0 38px 110px rgba(7,18,37,.22);backdrop-filter:blur(14px)}
    .login-brand{position:relative;padding:3rem 2.8rem;color:white;background:radial-gradient(circle at 18% 18%,rgba(245,180,51,.30),transparent 25%),linear-gradient(135deg,#06162f 0%,#082b59 52%,#0b4b91 100%);min-height:650px;display:flex;flex-direction:column;justify-content:space-between}
    .login-brand:before{content:"";position:absolute;inset:0;background:linear-gradient(120deg,rgba(255,255,255,.08) 0 1px,transparent 1px 18px),radial-gradient(circle at 86% 14%,rgba(255,255,255,.20),transparent 20%);opacity:.7;pointer-events:none}
    .brand-content,.brand-footer{position:relative;z-index:1}
    .login-logo-row{display:flex;align-items:center;gap:16px;margin-bottom:2rem}
    .login-logo-row img{width:86px;height:86px;border-radius:22px;background:white;padding:8px;object-fit:contain;box-shadow:0 18px 45px rgba(0,0,0,.28)}
    .login-kicker{font-size:.78rem;font-weight:900;text-transform:uppercase;letter-spacing:.18em;color:#f5b433;margin-bottom:.4rem}
    .login-brand h1{margin:0;font-size:2.65rem;line-height:1.04;letter-spacing:-.055em;color:white;font-weight:950}
    .login-brand p{font-size:1.03rem;line-height:1.65;color:#dbeafe;max-width:610px;margin:1.05rem 0}
    .login-badges{display:flex;gap:9px;flex-wrap:wrap;margin:1.25rem 0 0}
    .login-badge{display:inline-flex;align-items:center;border:1px solid rgba(255,255,255,.22);background:rgba(255,255,255,.12);border-radius:999px;padding:7px 11px;color:#fff;font-size:.78rem;font-weight:850}
    .login-feature-grid{display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:12px;margin-top:1.5rem}
    .login-feature{border:1px solid rgba(255,255,255,.16);background:rgba(255,255,255,.10);border-radius:18px;padding:13px 14px;color:#eaf2ff}
    .login-feature b{display:block;color:white;font-size:.92rem;margin-bottom:4px}.login-feature span{font-size:.78rem;color:#cfe1ff}
    .brand-footer{border-top:1px solid rgba(255,255,255,.18);padding-top:1rem;color:#cbd5e1;font-size:.82rem;display:flex;justify-content:space-between;gap:12px;flex-wrap:wrap}
    .login-panel{padding:3rem 2.6rem;background:linear-gradient(180deg,#ffffff 0%,#f8fbff 100%);display:flex;flex-direction:column;justify-content:center}
    .login-card{background:white;border:1px solid #dce6f2;border-radius:30px;padding:2rem;box-shadow:0 18px 55px rgba(15,23,42,.10)}
    .login-card h2{font-size:1.75rem;margin:0 0 .35rem;color:#071225;font-weight:950}.login-card .muted{color:#64748b;margin:0 0 1.25rem;line-height:1.55}
    .login-mini{display:grid;grid-template-columns:repeat(3,1fr);gap:10px;margin:1.1rem 0 0}.login-mini div{background:#f1f5f9;border:1px solid #dbe3ef;border-radius:16px;padding:10px;text-align:center}.login-mini b{display:block;color:#0b3b76;font-size:1rem}.login-mini span{font-size:.70rem;color:#64748b;font-weight:800;text-transform:uppercase;letter-spacing:.05em}
    .login-help{margin-top:1rem;padding:12px 14px;border-radius:18px;background:#fff8eb;border:1px solid #f3d79a;color:#6b4b0b;font-size:.86rem;line-height:1.5}
    .login-card div[data-testid="stForm"]{border:0;padding:0}.login-card label{font-weight:850;color:#0f172a!important}.login-card input{border-radius:14px!important}
    .login-card .stButton>button{width:100%;height:3rem;border-radius:16px;background:linear-gradient(135deg,#071225,#0b3b76);border:0;color:white;font-weight:950;letter-spacing:.02em;box-shadow:0 14px 32px rgba(11,59,118,.24)}
    .login-card .stButton>button:hover{background:linear-gradient(135deg,#04101f,#08315f);transform:translateY(-1px)}
    .login-demo{margin-top:1rem}.login-demo div[data-testid="stExpander"]{box-shadow:none;border-radius:18px;background:#f8fafc}
    @media(max-width:920px){.login-frame{grid-template-columns:1fr}.login-brand{min-height:auto;padding:2.1rem}.login-panel{padding:1.4rem}.login-brand h1{font-size:2rem}.login-feature-grid{grid-template-columns:1fr}.login-mini{grid-template-columns:1fr}}
    .stButton>button,.stDownloadButton>button{border-radius:13px;border:1px solid var(--psb-blue);background:var(--psb-blue);color:white;font-weight:800;box-shadow:0 8px 18px rgba(11,59,118,.16)}
    .stButton>button:hover,.stDownloadButton>button:hover{background:var(--psb-navy);color:white;border-color:var(--psb-navy)}
    div[data-testid="stDataFrame"]{border-radius:18px;overflow:hidden;border:1px solid var(--psb-line);box-shadow:0 10px 26px rgba(15,23,42,.05)}
    div[data-testid="stExpander"]{border-radius:18px;border:1px solid var(--psb-line);background:white;box-shadow:0 8px 22px rgba(15,23,42,.04)}
    .stTabs [data-baseweb="tab-list"]{gap:8px}
    .stTabs [data-baseweb="tab"]{border-radius:999px;background:#e8eef7;padding:.45rem 1rem;font-weight:800}
    h1,h2,h3{letter-spacing:-.025em;color:#0f172a}
    </style>
    """, unsafe_allow_html=True)


def header() -> None:
    logo = f"<img src='{logo_data_uri()}' />" if LOGO_PATH.exists() else ""
    st.markdown(f"""
    <div class='psb-hero'>{logo}<div>
    <h1>{APP_TITLE}</h1><p>{APP_SUBTITLE}</p>
    <div>{"".join([f"<span class='pill'>{s}</span>" for s in STANDARDS])}{backend_status_badges()}</div>
    </div></div>
    """, unsafe_allow_html=True)


def table(df: pd.DataFrame, max_rows: int = 300) -> None:
    if df is None or df.empty:
        st.caption("No records found.")
        return
    shown = df.fillna("")
    if len(shown) > max_rows:
        st.caption(f"Showing latest {max_rows} of {len(shown)} records for faster loading. Use Backup/Export for full data.")
        shown = shown.tail(max_rows)
    st.dataframe(shown, width="stretch", hide_index=True)


def metrics(items):
    cols = st.columns(4)
    for i, (label, value) in enumerate(items):
        cols[i % 4].metric(label, value)


def login_page() -> None:
    if "captcha_question" not in st.session_state:
        a, b = random.randint(2, 12), random.randint(2, 12)
        st.session_state["captcha_question"] = f"{a} + {b}"
        st.session_state["captcha_answer"] = str(a + b)

    logo_html = f"<img src='{logo_data_uri()}' alt='PSB Logo' />" if LOGO_PATH.exists() else ""
    standards_html = "".join([f"<span class='login-badge'>{s}</span>" for s in STANDARDS[:6]])

    st.markdown(f"""
    <div class='login-shell'>
      <div class='login-frame'>
        <section class='login-brand'>
          <div class='brand-content'>
            <div class='login-logo-row'>
              {logo_html}
              <div>
                <div class='login-kicker'>Maritime Training & Competency</div>
                <div style='font-weight:900;color:#fff;font-size:1.05rem'>Pakistan Shipping Bureau</div>
              </div>
            </div>
            <h1>Classification Society HRDM Platform</h1>
            <div class='login-badges'>{standards_html}</div>
          </div>
          <div class='brand-footer'>
            <span>Secure Role-Based Access</span><span>ISO / IACS Ready</span>
          </div>
        </section>
        <section class='login-panel'>
          <div class='login-card'>
            <h2>Sign In</h2>
            <p class='muted'>Access your account</p>
    """, unsafe_allow_html=True)

    login_attempts = st.session_state.get("login_attempts", 0)
    blocked_until = st.session_state.get("login_blocked_until")
    now_ts = datetime.utcnow()
    if blocked_until and isinstance(blocked_until, datetime) and now_ts < blocked_until:
        remaining = int((blocked_until - now_ts).total_seconds() / 60) + 1
        st.error(f"Too many failed login attempts. Please try again in {remaining} minute(s).")

    with st.form("login", clear_on_submit=False):
        login = st.text_input("Login ID or Email", placeholder="Enter your login ID or official email")
        password = st.text_input("Password", type="password", placeholder="Enter your password")
        captcha = st.text_input(f"Security Verification: {st.session_state['captcha_question']} = ?", placeholder="Answer")
        submit = st.form_submit_button("Sign in to PSB Portal")

    if submit:
        if blocked_until and isinstance(blocked_until, datetime) and now_ts < blocked_until:
            st.error("You are temporarily blocked due to too many failed login attempts.")
            return
        if captcha.strip() != st.session_state.get("captcha_answer", ""):
            st.error("Security verification failed. Please try again.")
            return
        login_value = clean(login).lower().strip()
        password_value = clean(password)
        if not login_value or not password_value:
            st.error("Login ID/email and password are required.")
            return
        match = db_where(
            "users",
            "(lower(login_id) = :login_key or lower(email) = :login_key) and password_hash = :password_hash and status = 'Active'",
            (("login_key", login_value), ("password_hash", phash(password_value))),
        )
        if match.empty:
            st.session_state["login_attempts"] = login_attempts + 1
            if st.session_state["login_attempts"] >= MAX_LOGIN_ATTEMPTS:
                st.session_state["login_blocked_until"] = now_ts + timedelta(minutes=LOGIN_BLOCK_MINUTES)
                st.error(f"Too many failed attempts. Try again after {LOGIN_BLOCK_MINUTES} minute(s).")
            else:
                remaining = MAX_LOGIN_ATTEMPTS - st.session_state["login_attempts"]
                st.error(f"Invalid login ID/email or password. {remaining} attempt(s) remaining.")
        else:
            user = match.iloc[0].to_dict()
            st.session_state["logged_in"] = True
            st.session_state["user"] = user
            st.session_state["login_attempts"] = 0
            st.session_state["login_blocked_until"] = None
            token = create_auth_token(user["user_id"])
            st.session_state["auth_token"] = token
            getattr(st, "experimental_set_query_params", lambda **kwargs: None)(auth_token=token)
            try:
                # set a secure cookie so auth survives simple page refreshes
                components.html(f"<script>document.cookie = 'psb_auth={token}; path=/; max-age=86400; Secure; SameSite=Strict';</script>", height=0)
            except Exception:
                pass
            db_update("users", "user_id", user["user_id"], {"last_login": now()})
            audit("User Login", f"{user['name']} logged in", actor=user)
            st.rerun()

    st.markdown("""
            <div class='login-mini'>
              <div><b>ISO</b><span>QMS</span></div>
              <div><b>IACS</b><span>Standards</span></div>
              <div><b>IMO</b><span>RO Code</span></div>
            </div>
          </div>
        </section>
      </div>
    </div>
    """, unsafe_allow_html=True)


def require_login() -> dict:
    if "logged_in" not in st.session_state:
        st.session_state["logged_in"] = False
    if "user" not in st.session_state:
        st.session_state["user"] = {}
    if not st.session_state["logged_in"]:
        params = getattr(st, "experimental_get_query_params", lambda: {})()
        token = clean(params.get("auth_token", [""])[0]) if params.get("auth_token") else ""
        # If no auth token in URL, try reading a browser cookie and reload with it (helps preserve login across refresh)
        if not token:
            try:
                components.html("""
                    <script>
                    (function(){
                        function getCookie(n){return document.cookie.split('; ').reduce(function(r,c){var p=c.split('='); return p[0]===n?decodeURIComponent(p.slice(1).join('=')):r},'');}
                        var t = getCookie('psb_auth');
                        if(t && !new URLSearchParams(window.location.search).has('auth_token')){
                            var params = new URLSearchParams(window.location.search);
                            params.set('auth_token', t);
                            window.location.search = params.toString();
                        }
                    })();
                    </script>
                """, height=0)
            except Exception:
                pass
        if token:
            user_id = resolve_auth_token(token)
            if user_id:
                match = db_where("users", "user_id = :user_id and status = 'Active'", (("user_id", user_id),))
                if not match.empty:
                    user = match.iloc[0].to_dict()
                    st.session_state["logged_in"] = True
                    st.session_state["user"] = user
        if not st.session_state["logged_in"]:
            # runtime checks for exposed secrets in environment (warn only)
            warnings = []
            if DEFAULT_ADMIN_PASSWORD:
                warnings.append("DEFAULT_ADMIN_PASSWORD is set — ensure this is not a weak/demo password and is provided via environment variables only.")
            if SUPABASE_SERVICE_ROLE_KEY:
                warnings.append("SUPABASE_SERVICE_ROLE_KEY is present in environment — keep this secret in hosting environment variables, not in code or repo.")
            if warnings:
                for w in warnings:
                    st.warning(w)
                st.info("Security: remove any hard-coded secrets and use environment variables. See .env.example and DEPLOYMENT_CHECKLIST.md for guidance.")
            login_page()
            st.stop()
    return st.session_state["user"]


def sidebar(actor: dict) -> str:
    if LOGO_PATH.exists():
        st.sidebar.image(str(LOGO_PATH), width=95)
    st.sidebar.success(f"{actor_get(actor,'name')} ({actor_get(actor,'role')})")
    st.sidebar.caption(actor_get(actor, "email"))
    role = actor_get(actor, "role")
    if role == "Admin":
        pages = ["Dashboard","Admin","Training Matrix","Training","Files","Development Plans","Competency","Practical/Witness","Authorization","CRB","Technical Authority","Survey Report Review","Plan Review QA","Competency NCR","Gap Advisor","Annual Board","Restrictions","Client Feedback","Succession","Workforce Planning","Accreditation Readiness","Interpretation Portal","Job Allocation","KPI","CPD","Knowledge Library","QMS","Revalidation","Backup","QR Verify","Management"]
    elif role == "Trainer":
        pages = ["Dashboard","Training Matrix","Training","Files","KPI","CPD","Knowledge Library"]
    elif role in ["Tutor/Mentor","Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Technical Manager"]:
        pages = ["Dashboard","Development Plans","Practical/Witness","Competency","Authorization","CRB","Technical Authority","Survey Report Review","Plan Review QA","Competency NCR","Gap Advisor","Annual Board","Interpretation Portal","Files","KPI","Knowledge Library"]
    elif role in ["QMR","QMS Auditor"]:
        pages = ["Dashboard","QMS","Authorization","CRB","Competency NCR","Gap Advisor","Annual Board","Accreditation Readiness","Knowledge Library","Revalidation","Backup"]
    elif role == "Job Coordinator":
        pages = ["Dashboard","Job Allocation","Workforce Planning","Client Feedback","KPI","Management"]
    elif role == "Management":
        pages = ["Dashboard","Management","Authorization","CRB","Technical Authority","Annual Board","Restrictions","Client Feedback","Succession","Workforce Planning","Accreditation Readiness","Job Allocation","KPI","Revalidation","Backup"]
    else:
        pages = ["Dashboard","Training","Files","Development Plans","Practical/Witness","Competency","Authorization","CPD","Knowledge Library"]
    page = st.sidebar.radio("Menu", pages)
    if st.sidebar.button("Logout"):
        audit("User Logout", f"{actor_get(actor,'name')} logged out", actor=actor)
        clear_auth_token()
        st.session_state["logged_in"] = False
        st.session_state["user"] = {}
        st.rerun()
    return page


def dashboard_page(actor):
    st.header(f"{actor_get(actor,'role')} Dashboard")
    users = db_all("users"); trainings = db_all("trainings"); records = db_all("training_records")
    comp = db_all("competency_matrix"); auths = db_all("authorization_requests"); jobs = db_all("job_requests")
    cpd = db_all("cpd_records"); kpi = db_all("kpi_records")
    notifications = db_all("notifications")
    metrics([
        ("Users", len(users)), ("Trainings", len(trainings)), ("Training Records", len(records)),
        ("Competencies", len(comp)), ("Approved Auth", len(auths[auths["status"]=="Management Approved"]) if not auths.empty else 0),
        ("Jobs Assigned", len(jobs[jobs["status"]=="Assigned"]) if not jobs.empty else 0),
        ("CPD Records", len(cpd)), ("KPI Records", len(kpi)),
    ])
    my_notifications = notifications[notifications["user_id"] == actor_get(actor, "user_id")] if not notifications.empty else pd.DataFrame()
    if not my_notifications.empty:
        st.subheader("My Notifications / Messages")
        show_cols = [c for c in ["created_on", "subject", "message", "type", "status"] if c in my_notifications.columns]
        table(my_notifications.sort_values("created_on", ascending=False).head(10)[show_cols])
    st.subheader("World-Class Qualification Flow")
    for i, s in enumerate([
        "Admin assigns role, path, mentor and authorization matrix.",
        "Trainer assigns theoretical training modules and assessments.",
        "Candidate passes all required theoretical modules.",
        "Candidate becomes eligible for witness surveys.",
        "Tutor records minimum witness surveys and performance.",
        "Candidate completes supervised survey or plan review exercises.",
        "Tutor / Principal / Technical Manager recommend to CRB.",
        "CRB reviews evidence and QMR validates QMS compliance.",
        "Management approves authorization and QR certificate is issued.",
        "Job Coordinator allocates work only by valid scope, level, KPI, risk and availability.",
        "Annual review, CPD, refresher training and reauthorization maintain competence.",
    ], 1):
        st.markdown(f"<div class='step'><b>{i}.</b> {s}</div>", unsafe_allow_html=True)


def admin_page(actor):
    st.header("Admin Control Center")
    st.subheader("Backend Persistence Status")
    c1, c2, c3 = st.columns(3)
    c1.metric("Database", "Persistent" if database_is_persistent() else "Local/Temporary")
    c2.metric("File Storage", "Persistent" if storage_is_persistent() else "Local/Missing")
    c3.metric("Runtime", "Render" if is_render_runtime() else "Local")
    if not database_is_persistent():
        st.warning("Local SQLite is only for testing. On Render, use Supabase/PostgreSQL DATABASE_URL to prevent data loss.")
    if not storage_is_persistent():
        st.warning("Supabase Storage is recommended for uploaded files. Local uploads may not persist on hosting platforms.")

    users = db_all("users")
    with st.form("create_user"):
        c1, c2 = st.columns(2)
        name = c1.text_input("Name")
        email = c2.text_input("Email")
        role = c1.selectbox("Role", ROLES)
        path = c2.selectbox("Trainee / Competency Path", [""] + TRAINEE_PATHS)
        dept = c1.text_input("Department", "Survey")
        duty = c2.text_input("Assigned Duty / Scope")
        mentors = users[users["role"].isin(["Trainer","Tutor/Mentor","Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Technical Manager"])] if not users.empty else pd.DataFrame()
        mentor = c1.selectbox("Assigned Mentor/Tutor", [""] + (mentors["name"].astype(str)+" — "+mentors["user_id"].astype(str)).tolist())
        location = c2.text_input("Location", "Karachi")
        password = st.text_input("Password blank=auto", type="password")
        submit = st.form_submit_button("Create User")
    if submit and name and email:
        # validate email and rate-limit admin user creation
        if not validate_email(email):
            st.error("Invalid email address.")
            return
        try:
            admin_key = actor_get(actor, "user_id") or "anon"
            if not RATE_LIMITER.allowed("create_user", admin_key, limit=5, per_seconds=60):
                st.error("Rate limit exceeded for creating users. Try again later.")
                return
        except Exception:
            pass
        login = re.sub(r"[^a-z0-9]", "", name.lower().replace(" ", ".")) or f"user{random.randint(100,999)}"
        password = password or temp_password()
        mentor_name, mentor_id = ("","")
        if mentor:
            mentor_name, mentor_id = mentor.split(" — ")
        db_insert("users", {
            "user_id": uid("USR"), "name": name, "role": role, "trainee_path": path, "department": dept,
            "assigned_duty": duty, "email": email, "login_id": login, "password_hash": phash(password),
            "temp_password": password, "status": "Active", "availability": "Available", "current_location": location,
            "mentor_id": mentor_id, "mentor_name": mentor_name, "competency_level": "Level 0 - Trainee",
            "created_on": today(), "last_login": "",
        })
        audit("User Created", f"{name} as {role} path {path}", actor=actor)
        st.success("User created.")
        st.code(f"Login: {login}\nPassword: {password}")
    st.subheader("Availability and Competency Level")
    users = db_all("users")
    if not users.empty:
        person = st.selectbox("Person", users["name"].astype(str)+" — "+users["user_id"].astype(str))
        pid = person.split(" — ")[-1]
        user = users[users["user_id"] == pid].iloc[0]
        c1, c2, c3 = st.columns(3)
        availability = c1.selectbox("Availability", ["Available","Busy","On Leave","Unavailable"], index=["Available","Busy","On Leave","Unavailable"].index(user["availability"]) if user["availability"] in ["Available","Busy","On Leave","Unavailable"] else 0)
        location = c2.text_input("Location", user["current_location"])
        level = c3.selectbox("Competency Level", COMPETENCY_LEVELS, index=COMPETENCY_LEVELS.index(user["competency_level"]) if user["competency_level"] in COMPETENCY_LEVELS else 0)
        if st.button("Update Person Status"):
            db_update("users", "user_id", pid, {"availability": availability, "current_location": location, "competency_level": level})
            audit("User Status Updated", pid, actor=actor)
            st.success("Updated.")
    table(db_all("users").drop(columns=["password_hash"], errors="ignore"))


def file_upload_panel(actor, linked_table="general", linked_id="general", category="Other"):
    cat = st.selectbox("File Category", FILE_CATEGORIES, index=FILE_CATEGORIES.index(category) if category in FILE_CATEGORIES else 0)
    uploads = st.file_uploader("Upload PDF, PPT/PPTX, DOC/DOCX, TXT, images, video or Excel", type=ALLOWED_EXTENSIONS, accept_multiple_files=True)
    if st.button("Upload File(s)"):
        if not uploads:
            st.error("Select file(s).")
        else:
            count = 0
            for f in uploads:
                try:
                    upload_file(f, actor, linked_table, linked_id, cat)
                    count += 1
                except Exception as e:
                    st.error(f"{f.name}: {e}")
            st.success(f"{count} file(s) uploaded.")


def files_page(actor):
    st.header("File Repository")
    linked_table = st.text_input("Linked Table", "general")
    linked_id = st.text_input("Linked ID", "general")
    file_upload_panel(actor, linked_table, linked_id, "Training Material")
    table(db_all("files"))


def training_matrix_page(actor):
    st.header("Theoretical Training Matrix")
    st.info("Admin/Trainer/Tutor/Mentor can add, edit, or delete theoretical modules. These modules make candidates eligible for witness survey only after passing.")
    role = actor_get(actor, "role")
    modules = db_all("training_modules")
    target_path_options = ["All"] + TRAINEE_PATHS + ["Surveyor", "Plan Appraiser", "Auditor", "Rule Development Rep", "Industrial Surveyor"]

    if role in ["Admin", "Trainer", "Tutor/Mentor"]:
        with st.expander("Add New Module"):
            with st.form("module_add"):
                c1, c2 = st.columns(2)
                title = c1.text_input("Module Title")
                group = c2.text_input("Module Group", "Technical")
                target_path = c1.selectbox("Target Path", target_path_options)
                custom_path = c1.text_input("Or custom target path")
                if clean(custom_path):
                    target_path = clean(custom_path)
                hours = c2.number_input("CPD Hours", 0.0, 100.0, 2.0)
                mandatory = c1.checkbox("Mandatory", True)
                refresher = c2.checkbox("Refresher Required", True)
                validity = c1.number_input("Validity Months", 1, 120, 36)
                submit = st.form_submit_button("Add Module")
            if submit and title:
                db_insert("training_modules", {
                    "module_id": uid("MOD"), "title": title, "module_group": group, "target_path": target_path,
                    "mandatory": "Yes" if mandatory else "No", "refresher_required": "Yes" if refresher else "No",
                    "cpd_hours": hours, "validity_months": validity, "added_by": actor_get(actor, "name"), "created_on": today(),
                })
                audit("Training Module Added", title, actor=actor)
                st.success("Module added.")

        if not modules.empty:
            st.subheader("Edit Existing Module")
            selected = st.selectbox("Select Module to Edit", modules["title"].astype(str) + " — " + modules["module_id"].astype(str))
            if selected:
                module_id = selected.split(" — ")[-1]
                module = modules[modules["module_id"] == module_id].iloc[0]
                default_index = target_path_options.index(module["target_path"]) if module["target_path"] in target_path_options else 0
                with st.form("module_edit"):
                    c1, c2 = st.columns(2)
                    title = c1.text_input("Module Title", module["title"])
                    group = c2.text_input("Module Group", module["module_group"])
                    selected_target = c1.selectbox("Target Path", target_path_options, index=default_index)
                    custom_path = c1.text_input("Or custom target path", "" if module["target_path"] in target_path_options else module["target_path"])
                    target_path = clean(custom_path) or selected_target
                    hours = c2.number_input("CPD Hours", 0.0, 100.0, float(module["cpd_hours"] or 0.0))
                    mandatory = c1.checkbox("Mandatory", module["mandatory"] == "Yes")
                    refresher = c2.checkbox("Refresher Required", module["refresher_required"] == "Yes")
                    validity = c1.number_input("Validity Months", 1, 120, int(module["validity_months"] or 36))
                    update = st.form_submit_button("Save Module Changes")
                if update:
                    db_update("training_modules", "module_id", module_id, {
                        "title": title, "module_group": group, "target_path": target_path,
                        "mandatory": "Yes" if mandatory else "No", "refresher_required": "Yes" if refresher else "No",
                        "cpd_hours": hours, "validity_months": validity, "updated_on": now(),
                    })
                    audit("Training Module Updated", title, actor=actor)
                    st.success("Module updated.")
                    st.rerun()
                if st.button("Delete Module", key="delete_module"):
                    db_delete("training_modules", "module_id", module_id)
                    audit("Training Module Deleted", module["title"], actor=actor)
                    st.success("Module deleted.")
                    st.rerun()

    table(db_all("training_modules"))


def training_page(actor):
    st.header("Training Management")
    role = actor_get(actor, "role")
    users = db_all("users"); trainings = db_all("trainings")
    if role in ["Admin","Trainer","Tutor/Mentor"]:
        with st.expander("Create Course from Theoretical Module"):
            modules = db_all("training_modules")
            trainers = users[(users["role"] == "Trainer") & (users["status"] == "Active")] if not users.empty else pd.DataFrame()
            with st.form("course"):
                module_sel = st.selectbox("Module", modules["title"].astype(str)+" — "+modules["module_id"].astype(str)) if not modules.empty else ""
                trainer = st.selectbox("Trainer", trainers["name"].astype(str)+" — "+trainers["user_id"].astype(str)) if not trainers.empty else ""
                target_roles = st.multiselect("Target Roles", ROLES, default=["Trainee"])
                passing = st.number_input("Passing Marks", 1, 100, 75)
                lms = st.text_input("LMS/SCORM Course ID")
                submit = st.form_submit_button("Create Course")
            if submit and module_sel and trainer:
                title, module_id = module_sel.split(" — ")
                trainer_name, trainer_id = trainer.split(" — ")
                module = modules[modules["module_id"] == module_id].iloc[0]
                tid = uid("TRN")
                db_insert("trainings", {
                    "training_id": tid, "module_id": module_id, "title": title, "category": module["module_group"],
                    "standards": join_list(STANDARDS), "target_roles": join_list(target_roles),
                    "target_paths": module["target_path"], "trainer_id": trainer_id, "trainer_name": trainer_name,
                    "slides_link": "", "video_link": "", "reference_link": "", "scorm_package_link": "",
                    "lms_course_id": lms, "schedule_date": "", "schedule_time": "10:00", "meeting_link": "",
                    "recording_link": "", "passing_marks": passing, "validity_months": module["validity_months"],
                    "max_attempts": 3, "retest_wait_days": 7, "status": "Draft", "created_on": now(), "updated_on": now(),
                })
                audit("Training Created", title, actor=actor)
                st.success("Course created.")
    trainings = db_all("trainings")
    if trainings.empty:
        st.warning("No training created.")
        return
    if role == "Trainer":
        trainings = trainings[trainings["trainer_id"] == actor_get(actor, "user_id")]
    elif role not in ["Admin","Trainer","Tutor/Mentor"]:
        rec = db_all("training_records")
        ids = rec[rec["user_id"] == actor_get(actor, "user_id")]["training_id"].tolist() if not rec.empty else []
        trainings = trainings[trainings["training_id"].isin(ids)]
    if trainings.empty:
        st.warning("No training assigned.")
        return
    selected = st.selectbox("Select Training", trainings["title"].astype(str)+" — "+trainings["training_id"].astype(str))
    tid = selected.split(" — ")[-1]
    tr = db_all("trainings")
    tr_row = tr[tr["training_id"] == tid].iloc[0]
    if role in ["Admin","Trainer","Tutor/Mentor"]:
        st.subheader("Edit Training Details")
        trainers = users[(users["role"] == "Trainer") & (users["status"] == "Active")] if not users.empty else pd.DataFrame()
        trainer_options = list(trainers["name"].astype(str) + " — " + trainers["user_id"].astype(str)) if not trainers.empty else [f"{tr_row['trainer_name']} — {tr_row['trainer_id']}"]
        trainer_default = f"{tr_row['trainer_name']} — {tr_row['trainer_id']}"
        trainer_index = trainer_options.index(trainer_default) if trainer_default in trainer_options else 0
        status_options = ["Draft", "Scheduled", "Completed", "Cancelled"]
        status_default = tr_row["status"] if tr_row["status"] in status_options else "Draft"
        status_index = status_options.index(status_default)
        with st.expander("Edit Course Details", expanded=False):
            c1, c2 = st.columns(2)
            title = c1.text_input("Training Title", tr_row["title"])
            category = c2.text_input("Category", tr_row["category"])
            trainer_selected = c1.selectbox("Trainer", trainer_options, index=trainer_index)
            target_roles = c2.multiselect("Target Roles", ROLES, default=split_list(tr_row["target_roles"]))
            target_paths = c1.text_input("Target Paths", tr_row["target_paths"])
            passing = c2.number_input("Passing Marks", 1, 100, int(tr_row["passing_marks"] or 75))
            validity = c1.number_input("Validity Months", 1, 120, int(tr_row["validity_months"] or 36))
            status = c2.selectbox("Status", status_options, index=status_index)
            if st.button("Save Training Details", key="save_training_details"):
                trainer_name, trainer_id = trainer_selected.split(" — ")
                db_update("trainings", "training_id", tid, {
                    "title": title, "category": category, "target_roles": join_list(target_roles),
                    "target_paths": target_paths, "trainer_id": trainer_id, "trainer_name": trainer_name,
                    "passing_marks": passing, "validity_months": validity, "status": status,
                    "updated_on": now(),
                })
                audit("Training Updated", title, actor=actor)
                st.success("Training details saved.")
            if st.button("Delete Training", key="delete_training"):
                db_delete("trainings", "training_id", tid)
                audit("Training Deleted", tr_row["title"], actor=actor)
                st.success("Training deleted.")
                st.rerun()
    if role in ["Admin","Trainer"]:
        tabs = st.tabs(["Files & Links","MCQ","Assignment","Attendance/Records"])
        with tabs[0]:
            file_upload_panel(actor, "trainings", tid, "Training Material")
            slides = st.text_input("Slides Link", tr_row["slides_link"])
            video = st.text_input("Video Link", tr_row["video_link"])
            ref = st.text_input("Reference Link", tr_row["reference_link"])
            scorm = st.text_input("SCORM Package Link", tr_row["scorm_package_link"])
            sdate = st.date_input("Schedule Date")
            stime = st.text_input("Schedule Time", tr_row["schedule_time"])
            st.link_button("Open MS Teams to Create Meeting", f"https://teams.microsoft.com/l/meeting/new?subject={quote_plus(clean(tr_row['title']))}")
            meeting = st.text_input("Final MS Teams Meeting Link", tr_row["meeting_link"])
            recording = st.text_input("Recording Link", tr_row["recording_link"])
            if st.button("Save Links and Schedule"):
                db_update("trainings", "training_id", tid, {"slides_link": slides, "video_link": video, "reference_link": ref, "scorm_package_link": scorm, "schedule_date": str(sdate), "schedule_time": stime, "meeting_link": meeting, "recording_link": recording, "status": "Scheduled", "updated_on": now()})
                st.success("Saved.")
            f = db_all("files")
            table(f[f["linked_id"] == tid] if not f.empty else f)
        with tabs[1]:
            st.info("Upload training source files here to extract text and generate MCQs for the selected course.")
            uploads = st.file_uploader("Upload source files for MCQ generation", type=ALLOWED_EXTENSIONS, accept_multiple_files=True, key="mcq_source_files")
            if st.button("Upload MCQ Source File(s)"):
                if not uploads:
                    st.error("Select file(s) to upload.")
                else:
                    uploaded = 0
                    for f in uploads:
                        try:
                            upload_file(f, actor, "trainings", tid, "Training Material")
                            uploaded += 1
                        except Exception as e:
                            st.error(f"{f.name}: {e}")
                    st.success(f"{uploaded} source file(s) uploaded.")
                    st.rerun()
            f = db_all("files")
            extracted = "\n".join(f[(f["linked_id"] == tid) & (f["extracted_text"] != "")]["extracted_text"].astype(str).tolist()) if not f.empty else ""
            content = st.text_area("MCQ Content", value=extracted, height=220)
            count = st.slider("Number of MCQs", 5, 30, 10)
            if st.button("Generate MCQs"):
                qs = generate_mcqs(tid, content, count)
                if qs.empty:
                    st.error("Could not generate MCQs. Add clearer text or upload a clearer source file.")
                else:
                    exec_sql("delete from question_bank where training_id=:tid", {"tid": tid})
                    for _, q in qs.iterrows():
                        db_insert("question_bank", q.to_dict())
                    st.success(f"{len(qs)} MCQs generated.")
            q = db_all("question_bank")
            training_qs = q[q["training_id"] == tid] if not q.empty else pd.DataFrame()
            if training_qs.empty:
                st.warning("No MCQs generated yet for this training.")
            else:
                st.subheader("Generated MCQs")
                table(training_qs)
                selected_question = st.selectbox("Select MCQ to delete", training_qs["question"].astype(str) + " — " + training_qs["question_id"].astype(str))
                if st.button("Delete Selected MCQ"):
                    qid = selected_question.split(" — ")[-1]
                    db_delete("question_bank", "question_id", qid)
                    st.success("MCQ deleted.")
                    st.rerun()
                st.markdown("---")
                st.subheader("Broadcast MCQs")
                recipient_roles = st.multiselect("Recipient Roles", ROLES, default=["Trainee"])
                recipients = users[(users["status"] == "Active") & (users["role"].isin(recipient_roles))] if not users.empty else pd.DataFrame()
                selected_receivers = st.multiselect("Send To", recipients["name"].astype(str) + " — " + recipients["user_id"].astype(str))
                broadcast_msg = st.text_area("Broadcast Message", f"New MCQs generated for {tr_row['title']}. Please login to review the course and attempt the assessment.")
                if st.button("Broadcast MCQs"):
                    if not selected_receivers:
                        st.error("Select at least one recipient.")
                    else:
                        sent = 0
                        for item in selected_receivers:
                            name, uidv = item.split(" — ")
                            create_notification(uidv, f"New MCQs Available: {tr_row['title']}", broadcast_msg, "MCQ Broadcast")
                            sent += 1
                        st.success(f"MCQs broadcast sent to {sent} recipients.")
        with tabs[2]:
            eligible = users[(users["status"] == "Active") & (users["role"].isin(split_list(tr_row["target_roles"])))] if not users.empty else pd.DataFrame()
            st.caption("You can assign by role/person. Admin/Trainer may add multiple theoretical modules before witness eligibility.")
            selected_users = st.multiselect("Assign Persons", eligible["name"].astype(str)+" — "+eligible["user_id"].astype(str))
            due = st.date_input("Due Date", date.today()+timedelta(days=30))
            if st.button("Assign Training"):
                records = db_all("training_records")
                added = 0
                for item in selected_users:
                    name, uidv = item.split(" — ")
                    if not records.empty and not records[(records["user_id"] == uidv) & (records["training_id"] == tid)].empty:
                        continue
                    u = users[users["user_id"] == uidv].iloc[0]
                    db_insert("training_records", {
                        "record_id": uid("REC"), "user_id": uidv, "name": name, "role": u["role"],
                        "trainee_path": u["trainee_path"], "training_id": tid, "training_title": tr_row["title"],
                        "status": "Pending", "slides_opened": "No", "video_opened": "No", "live_attendance": "Not Marked",
                        "recording_opened": "No", "lms_completed": "No", "test_status": "Not Attempted", "score": None,
                        "passing_marks": tr_row["passing_marks"], "certificate_status": "Not Issued", "certificate_link": "",
                        "due_date": str(due), "completed_on": "", "progress": 0, "remarks": "Assigned", "updated_on": now(),
                    })
                    assignment_msg = (
                        f"You have been assigned training: {tr_row['title']}. "
                        f"Schedule: {clean(tr_row.get('schedule_date')) or 'To be announced'} "
                        f"at {clean(tr_row.get('schedule_time')) or 'To be announced'}. "
                        f"Due date: {due}. Please open the Training Management page to access read-only materials, meeting link, recording link and assessment."
                    )
                    create_notification(uidv, f"Training Assigned: {tr_row['title']}", assignment_msg, "Training")
                    added += 1
                st.success(f"{added} persons assigned.")
        with tabs[3]:
            rec = db_all("training_records")
            assigned = rec[rec["training_id"] == tid] if not rec.empty else pd.DataFrame()
            table(assigned)
            if not assigned.empty:
                person = st.selectbox("Mark Attendance", assigned["name"].astype(str)+" — "+assigned["user_id"].astype(str))
                att = st.selectbox("Attendance", ["Present", "Absent"])
                if st.button("Save Attendance"):
                    uidv = person.split(" — ")[-1]
                    rr = assigned[assigned["user_id"] == uidv].iloc[0]
                    db_update("training_records", "record_id", rr["record_id"], {"live_attendance": att, "updated_on": now()})
                    update_training_progress(rr["record_id"])
                    st.success("Attendance saved.")
    else:
        trainee_training(actor, tid)


def trainee_training(actor, tid):
    """Read-only trainee view.
    Trainees can see assigned training schedule and materials, but cannot edit course data.
    Opening/confirming material updates only their own training record.
    """
    uidv = actor_get(actor, "user_id")
    rr = db_where("training_records", "user_id = :user_id and training_id = :training_id", (("user_id", uidv), ("training_id", tid)))
    if rr.empty:
        st.warning("Training not assigned.")
        return
    tr = db_where("trainings", "training_id = :training_id", (("training_id", tid),))
    if tr.empty:
        st.warning("Training details not found.")
        return

    row = rr.iloc[0]
    tr_row = tr.iloc[0]
    record_id = row["record_id"]
    is_absent = clean(row.get("live_attendance")) == "Absent"

    st.subheader(clean(tr_row["title"]))
    metrics([
        ("Progress", f"{row['progress']}%"),
        ("Attendance", clean(row.get("live_attendance", "Not Marked"))),
        ("LMS", row["lms_completed"]),
        ("Test", row["test_status"]),
    ])

    st.info(
        f"Schedule: {clean(tr_row.get('schedule_date')) or 'Not scheduled'} "
        f"at {clean(tr_row.get('schedule_time')) or 'Not specified'} | "
        f"Trainer: {clean(tr_row.get('trainer_name')) or 'Not assigned'} | Due: {clean(row.get('due_date'))}"
    )

    st.markdown("### Live Session")
    meeting_link = clean(tr_row.get("meeting_link"))
    if meeting_link:
        st.link_button("Join / Open Meeting Link", meeting_link)
    else:
        st.caption("Meeting link is not available yet.")

    st.markdown("### Training Material (Read Only)")
    c1, c2, c3, c4 = st.columns(4)
    slides_link = clean(tr_row.get("slides_link"))
    video_link = clean(tr_row.get("video_link"))
    reference_link = clean(tr_row.get("reference_link"))
    scorm_link = clean(tr_row.get("scorm_package_link"))

    if slides_link:
        c1.link_button("Open Slides", slides_link)
        if c1.button("Confirm Slides Completed", key=f"slides_done_{record_id}"):
            db_update("training_records", "record_id", record_id, {"slides_opened": "Yes", "updated_on": now()})
            update_training_progress(record_id); st.rerun()
    else:
        c1.caption("Slides not uploaded.")

    if video_link:
        c2.link_button("Open Video", video_link)
        if c2.button("Confirm Video Completed", key=f"video_done_{record_id}"):
            db_update("training_records", "record_id", record_id, {"video_opened": "Yes", "updated_on": now()})
            update_training_progress(record_id); st.rerun()
    else:
        c2.caption("Video not uploaded.")

    if reference_link:
        c3.link_button("Open Reference", reference_link)
    else:
        c3.caption("Reference link not uploaded.")

    if scorm_link:
        c4.link_button("Open LMS/SCORM", scorm_link)
        if c4.button("Confirm LMS Completed", key=f"lms_done_{record_id}"):
            db_update("training_records", "record_id", record_id, {"lms_completed": "Yes", "updated_on": now()})
            update_training_progress(record_id); st.rerun()
    else:
        c4.caption("LMS/SCORM link not uploaded.")

    linked_files = db_where("files", "linked_table = :linked_table and linked_id = :linked_id", (("linked_table", "trainings"), ("linked_id", tid)))
    if not linked_files.empty:
        st.markdown("#### Uploaded Documents / Files")
        for _, f in linked_files.iterrows():
            file_url = clean(f.get("public_url"))
            file_name = clean(f.get("file_name"))
            if file_url:
                st.link_button(f"Open {file_name}", file_url)
            else:
                st.caption(file_name)

    st.markdown("### Recording for Absent / Revision")
    recording_link = clean(tr_row.get("recording_link"))
    if recording_link:
        st.link_button("Open Recording", recording_link)
        if st.button("Confirm Recording Viewed", key=f"recording_done_{record_id}"):
            patch = {"recording_opened": "Yes", "video_opened": "Yes", "updated_on": now()}
            if is_absent or clean(row.get("live_attendance")) in ["Not Marked", ""]:
                patch["live_attendance"] = "Recording Viewed"
            db_update("training_records", "record_id", record_id, patch)
            update_training_progress(record_id); st.rerun()
    elif is_absent:
        st.warning("You were marked absent. Recording will appear here after the trainer uploads/pastes the recording link.")
    else:
        st.caption("Recording link is not available yet.")

    st.markdown("### Assessment")
    qs = db_where("question_bank", "training_id = :training_id", (("training_id", tid),))
    if qs.empty:
        st.warning("MCQs not generated yet.")
        return
    if row["test_status"] == "Passed":
        st.success("Assessment already passed.")
        return

    history = db_where("assessment_history", "user_id = :user_id and training_id = :training_id", (("user_id", uidv), ("training_id", tid)))
    attempts = len(history) if not history.empty else 0
    with st.form(f"assessment_{tid}_{uidv}"):
        answers = {}
        for i, (_, q) in enumerate(qs.iterrows(), 1):
            st.markdown(f"**Q{i}. {q['question']}**")
            opts = [q["option_a"], q["option_b"], q["option_c"], q["option_d"]]
            answers[q["question_id"]] = st.radio("Select", opts, key=f"{tid}_{uidv}_{q['question_id']}", label_visibility="collapsed")
        submit = st.form_submit_button("Submit Assessment")
    if submit:
        correct = sum(1 for _, q in qs.iterrows() if answers.get(q["question_id"]) == q["correct_answer"])
        score = round(correct / len(qs) * 100, 2)
        result = "Passed" if score >= int(tr_row["passing_marks"]) else "Failed"
        db_insert("assessment_history", {"assessment_id": uid("ASM"), "user_id": uidv, "name": actor_get(actor,"name"), "training_id": tid, "training_title": tr_row["title"], "attempt_no": attempts+1, "score": score, "result": result, "attempted_on": now(), "next_retest_allowed": str(date.today()+timedelta(days=7)) if result=="Failed" else "", "remarks": f"Correct {correct}/{len(qs)}"})
        db_update("training_records","record_id",record_id,{"score":score,"test_status":result,"certificate_status":"Issued" if result=="Passed" else "Not Issued","certificate_link":f"{PUBLIC_URL}/training-certificates/{uidv}/{tid}" if result=="Passed" else "","remarks":f"Correct {correct}/{len(qs)}", "updated_on": now()})
        update_training_progress(record_id)
        st.success(f"{result}: {score}%")
        st.rerun()


def development_plan_page(actor):
    st.header("Development Plans and Field Exposure Matrix")
    users = db_all("users")
    allowed = actor_get(actor, "role") in ["Admin","Tutor/Mentor","Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Technical Manager"]
    if allowed:
        with st.form("plan"):
            candidates = users[users["role"].isin(["Trainee","On Probation","Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor","Rule Development Rep"])] if not users.empty else pd.DataFrame()
            person = st.selectbox("Person", candidates["name"].astype(str)+" — "+candidates["user_id"].astype(str)) if not candidates.empty else ""
            scope = st.selectbox("Scope", SCOPES)
            month_no = st.number_input("Month No.", 1, 24, 1)
            activity = st.text_area("Development Activity")
            target = st.date_input("Target Date", date.today()+timedelta(days=30))
            submit = st.form_submit_button("Add Development Plan Item")
        if submit and person:
            name, uidv = person.split(" — ")
            u = users[users["user_id"] == uidv].iloc[0]
            db_insert("development_plans", {"plan_id": uid("PLAN"), "user_id": uidv, "name": name, "trainee_path": u["trainee_path"], "mentor_id": actor_get(actor,"user_id"), "mentor_name": actor_get(actor,"name"), "competency_scope": scope, "month_no": month_no, "activity": activity, "target_date": str(target), "status": "Open", "mentor_comments": "", "created_on": now(), "updated_on": now()})
            st.success("Development plan item added.")
        with st.form("exposure"):
            person2 = st.selectbox("Exposure Person", candidates["name"].astype(str)+" — "+candidates["user_id"].astype(str), key="expperson") if not candidates.empty else ""
            scope2 = st.selectbox("Exposure Scope", SCOPES)
            activity_type = st.selectbox("Activity Type", ["Witness Survey","Supervised Survey","Joint Plan Review","Independent Plan Review","Witness Audit","Independent Audit","Rule Exercise"])
            required = st.number_input("Required Count", 0, 20, 2)
            submit2 = st.form_submit_button("Add/Update Exposure Requirement")
        if submit2 and person2:
            name, uidv = person2.split(" — ")
            u = users[users["user_id"] == uidv].iloc[0]
            db_insert("field_exposure_matrix", {"exposure_id": uid("EXP"), "user_id": uidv, "name": name, "trainee_path": u["trainee_path"], "scope": scope2, "activity_type": activity_type, "required_count": required, "completed_count": 0, "status": "Pending", "updated_on": now()})
            st.success("Exposure requirement added.")
    plans = db_all("development_plans")
    exposures = db_all("field_exposure_matrix")
    if not allowed:
        plans = plans[plans["user_id"] == actor_get(actor, "user_id")] if not plans.empty else plans
        exposures = exposures[exposures["user_id"] == actor_get(actor, "user_id")] if not exposures.empty else exposures
    st.subheader("Development Plans")
    table(plans)
    st.subheader("Field Exposure Matrix")
    table(exposures)


def competency_page(actor):
    st.header("Competency Matrix and Authorization Matrix")
    users = db_all("users")
    if actor_get(actor, "role") in ["Admin","Tutor/Mentor","Technical Manager","Principal Surveyor","Chief Plan Appraiser","QMR","Management"]:
        with st.form("competency"):
            eligible = users[users["role"].isin(["Trainee","On Probation","Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor","Rule Development Rep"])] if not users.empty else pd.DataFrame()
            person = st.selectbox("Person", eligible["name"].astype(str)+" — "+eligible["user_id"].astype(str)) if not eligible.empty else ""
            scope = st.selectbox("Scope", SCOPES)
            matrix = get_matrix_for_scope(scope)
            job_type = matrix["job_type"] if matrix is not None else st.selectbox("Job Type", JOB_TYPES)
            level = st.selectbox("Current Competency Level", COMPETENCY_LEVELS)
            expiry = st.date_input("Expiry Target", date.today()+timedelta(days=365*3))
            submit = st.form_submit_button("Add Competency Record")
        if submit and person:
            name, uidv = person.split(" — ")
            u = users[users["user_id"] == uidv].iloc[0]
            matrix = get_matrix_for_scope(scope)
            db_insert("competency_matrix", {
                "competency_id": uid("COMP"), "user_id": uidv, "name": name, "role": u["role"], "trainee_path": u["trainee_path"],
                "area": scope, "competency_level": level, "scope": scope, "job_type": job_type,
                "required_training_ids": "", "required_witness_count": int(matrix["required_witness_count"]) if matrix is not None else 2,
                "required_supervised_count": int(matrix["required_supervised_count"]) if matrix is not None else 1,
                "required_joint_plan_count": int(matrix["required_joint_plan_count"]) if matrix is not None else 0,
                "required_independent_plan_count": int(matrix["required_independent_plan_count"]) if matrix is not None else 0,
                "required_level_for_auth": matrix["required_level_for_auth"] if matrix is not None else "Level 3 - Authorized",
                "status": "Pending", "expiry_date": str(expiry), "evidence": "", "created_on": now(), "updated_on": now(),
            })
            st.success("Competency added.")
    comp = db_all("competency_matrix")
    if actor_get(actor, "role") not in ["Admin","Tutor/Mentor","Technical Manager","Principal Surveyor","Chief Plan Appraiser","QMR","Management"]:
        comp = comp[comp["user_id"] == actor_get(actor, "user_id")] if not comp.empty else comp
    table(comp)
    st.subheader("Scope-Specific Authorization Matrix")
    table(db_all("authorization_matrix"))


def practical_page(actor):
    st.header("Practical / Witness / Supervised Assessment")
    users = db_all("users")
    allowed = actor_get(actor, "role") in ["Admin","Tutor/Mentor","Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Technical Manager","Trainer"]
    tabs = st.tabs(["Witness Survey","Supervised / Plan Review","Readiness"])
    with tabs[0]:
        if allowed:
            with st.form("witness"):
                candidates = users[users["role"].isin(["Trainee","On Probation","Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor","Rule Development Rep"])] if not users.empty else pd.DataFrame()
                person = st.selectbox("Person", candidates["name"].astype(str)+" — "+candidates["user_id"].astype(str)) if not candidates.empty else ""
                vessel = st.text_input("Vessel / Project")
                job_type = st.selectbox("Job Type", JOB_TYPES)
                scope = st.selectbox("Scope", SCOPES)
                location = st.text_input("Location")
                tech = st.slider("Technical Knowledge", 1, 5, 3)
                rule = st.slider("Rule Application", 1, 5, 3)
                safety = st.slider("Safety Awareness", 1, 5, 3)
                comm = st.slider("Communication", 1, 5, 3)
                report = st.slider("Report Quality", 1, 5, 3)
                conduct = st.slider("Professional Conduct", 1, 5, 3)
                outcome = st.selectbox("Outcome", ["Pass","Conditional","Fail"])
                comments = st.text_area("Tutor Comments")
                submit = st.form_submit_button("Submit Witness Assessment")
            if submit and person:
                name, uidv = person.split(" — ")
                u = users[users["user_id"] == uidv].iloc[0]
                db_insert("witness_surveys", {"witness_id": uid("WIT"), "user_id": uidv, "name": name, "trainee_path": u["trainee_path"], "tutor_id": actor_get(actor,"user_id"), "tutor_name": actor_get(actor,"name"), "vessel_or_project": vessel, "job_type": job_type, "scope": scope, "witness_date": today(), "location": location, "technical_knowledge": tech, "rule_application": rule, "safety_awareness": safety, "communication": comm, "report_quality": report, "professional_conduct": conduct, "outcome": outcome, "comments": comments, "status": "Submitted", "created_on": now(), "updated_on": now()})
                st.success("Witness survey recorded.")
        w = db_all("witness_surveys")
        if not allowed:
            w = w[w["user_id"] == actor_get(actor,"user_id")] if not w.empty else w
        table(w)
    with tabs[1]:
        if allowed:
            with st.form("supervised"):
                candidates = users[users["role"].isin(["Trainee","On Probation","Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor","Rule Development Rep"])] if not users.empty else pd.DataFrame()
                person = st.selectbox("Person", candidates["name"].astype(str)+" — "+candidates["user_id"].astype(str), key="sup_person") if not candidates.empty else ""
                kind = st.selectbox("Activity Kind", ["Supervised Survey","Joint Plan Review","Independent Plan Review","Independent Audit","Supervised Rule Exercise"])
                project = st.text_input("Vessel / Plan / Audit / Project")
                job_type = st.selectbox("Job Type", JOB_TYPES, key="sup_job")
                scope = st.selectbox("Scope", SCOPES, key="sup_scope")
                location = st.text_input("Location", key="sup_loc")
                prep = st.slider("Preparation", 1, 5, 3)
                exe = st.slider("Execution Quality", 1, 5, 3)
                find = st.slider("Findings Quality", 1, 5, 3)
                rep = st.slider("Reporting Quality", 1, 5, 3)
                rule = st.slider("Rule Compliance", 1, 5, 3)
                outcome = st.selectbox("Outcome", ["Pass","Conditional","Fail"], key="sup_out")
                comments = st.text_area("Comments", key="sup_com")
                submit = st.form_submit_button("Submit Supervised Assessment")
            if submit and person:
                name, uidv = person.split(" — ")
                u = users[users["user_id"] == uidv].iloc[0]
                db_insert("supervised_activities", {"supervised_id": uid("SUP"), "user_id": uidv, "name": name, "trainee_path": u["trainee_path"], "tutor_id": actor_get(actor,"user_id"), "tutor_name": actor_get(actor,"name"), "activity_kind": kind, "vessel_or_project": project, "job_type": job_type, "scope": scope, "activity_date": today(), "location": location, "preparation": prep, "execution_quality": exe, "findings_quality": find, "reporting_quality": rep, "rule_compliance": rule, "outcome": outcome, "comments": comments, "status": "Submitted", "created_on": now(), "updated_on": now()})
                st.success("Supervised activity recorded.")
        sup = db_all("supervised_activities")
        if not allowed:
            sup = sup[sup["user_id"] == actor_get(actor,"user_id")] if not sup.empty else sup
        table(sup)
    with tabs[2]:
        users2 = users if allowed else users[users["user_id"] == actor_get(actor, "user_id")]
        if not users2.empty:
            person = st.selectbox("Check Person", users2["name"].astype(str)+" — "+users2["user_id"].astype(str))
            scope = st.selectbox("Readiness Scope", SCOPES, key="ready_scope")
            uidv = person.split(" — ")[-1]
            ok, gaps = readiness(uidv, scope)
            st.subheader("Readiness Result")
            if ok:
                st.success("READY FOR CRB / AUTHORIZATION")
            else:
                st.error("NOT READY")
                for g in gaps:
                    st.write("- " + g)


def authorization_page(actor):
    st.header("Authorization Workflow")
    comp = db_all("competency_matrix")
    if comp.empty:
        st.warning("No competency records.")
        return
    if actor_get(actor, "role") in ["Admin","Tutor/Mentor","Surveyor","Plan Appraiser","Trainee","On Probation","Technical Manager"]:
        eligible = comp if actor_get(actor, "role") in ["Admin","Tutor/Mentor","Technical Manager"] else comp[comp["user_id"] == actor_get(actor, "user_id")]
        if not eligible.empty:
            sel = st.selectbox("Competency", eligible["name"].astype(str)+" — "+eligible["scope"].astype(str)+" — "+eligible["competency_id"].astype(str))
            cid = sel.split(" — ")[-1]
            c = comp[comp["competency_id"] == cid].iloc[0]
            ok, gaps = readiness(c["user_id"], c["scope"])
            if ok:
                st.success("Evidence complete. Eligible for authorization request.")
            else:
                st.warning("Evidence gaps:")
                for g in gaps:
                    st.write("- " + g)
            if st.button("Create Authorization Request"):
                if not ok:
                    st.error("Cannot create authorization request until required training, witness, supervised/plan-review and development plan evidence are complete.")
                else:
                    matrix = get_matrix_for_scope(c["scope"])
                    expiry = add_months(int(matrix["validity_months"])) if matrix is not None else add_months(36)
                    db_insert("authorization_requests", {"authorization_id": uid("AUTH"), "user_id": c["user_id"], "name": c["name"], "trainee_path": c["trainee_path"], "job_type": c["job_type"], "scope": c["scope"], "competency_id": cid, "status": "Tutor Recommended", "tutor_remarks": "Submitted based on completed evidence.", "tutor_signature": actor_get(actor,"name"), "tutor_signed_on": now(), "principal_remarks": "", "principal_signature": "", "principal_signed_on": "", "technical_remarks": "", "technical_signature": "", "technical_signed_on": "", "qms_remarks": "", "qms_signature": "", "qms_signed_on": "", "crb_decision": "", "crb_remarks": "", "management_remarks": "", "management_signature": "", "management_signed_on": "", "expiry_date": expiry, "certificate_id": "", "certificate_html": "", "certificate_storage_link": "", "qr_data_uri": "", "created_on": now(), "updated_on": now()})
                    st.success("Authorization request created and tutor recommendation recorded.")
    auths = db_all("authorization_requests")
    table(auths)
    if auths.empty:
        return
    sel = st.selectbox("Select Request", auths["name"].astype(str)+" — "+auths["scope"].astype(str)+" — "+auths["authorization_id"].astype(str))
    aid = sel.split(" — ")[-1]
    req = auths[auths["authorization_id"] == aid].iloc[0]
    role = actor_get(actor, "role")
    current = req["status"]
    next_status = None; remarks_field = None; sig_field = None; signed_field = None
    if role in ["Principal Surveyor","Chief Plan Appraiser","Lead Auditor"] and current == "Tutor Recommended":
        next_status, remarks_field, sig_field, signed_field = "Principal Reviewed", "principal_remarks", "principal_signature", "principal_signed_on"
    elif role == "Technical Manager" and current in ["Tutor Recommended","Principal Reviewed"]:
        next_status, remarks_field, sig_field, signed_field = "Technical Reviewed", "technical_remarks", "technical_signature", "technical_signed_on"
    elif role == "QMR" and current == "Technical Reviewed":
        next_status, remarks_field, sig_field, signed_field = "QMS Reviewed", "qms_remarks", "qms_signature", "qms_signed_on"
    elif role in ["Management","Admin"] and current in ["CRB Approved","QMS Reviewed"]:
        next_status, remarks_field, sig_field, signed_field = "Management Approved", "management_remarks", "management_signature", "management_signed_on"
    remarks = st.text_area("Approval Remarks")
    signature = st.text_input("Digital Signature", actor_get(actor, "name"))
    if st.button("Approve Next Step"):
        if not next_status:
            st.error("Your role cannot approve the current stage.")
        else:
            patch = {"status": next_status, "updated_on": now(), remarks_field: remarks, sig_field: signature, signed_field: now()}
            if next_status == "Management Approved":
                tmp = req.copy()
                for k,v in patch.items():
                    tmp[k] = v
                cert_id, html, qr = build_certificate(tmp)
                patch.update({"certificate_id": cert_id, "certificate_html": html, "certificate_storage_link": f"database://authorization_certificates/{cert_id}", "qr_data_uri": qr})
                db_insert("authorization_certificates", {"certificate_id": cert_id, "authorization_id": aid, "user_id": req["user_id"], "name": req["name"], "scope": req["scope"], "job_type": req["job_type"], "issue_date": today(), "expiry_date": req["expiry_date"], "certificate_html": html, "qr_data_uri": qr, "storage_link": f"database://authorization_certificates/{cert_id}", "verification_url": f"{PUBLIC_URL}/verify/{cert_id}", "status": "Valid", "created_on": now()})
                db_update("competency_matrix","competency_id",req["competency_id"],{"status":"Authorized","competency_level":"Level 3 - Authorized","updated_on":now()})
                db_update("users","user_id",req["user_id"],{"competency_level":"Level 3 - Authorized"})
            db_update("authorization_requests","authorization_id",aid,patch)
            st.success(f"Moved to {next_status}")
            st.rerun()
    req2 = db_all("authorization_requests")
    req2 = req2[req2["authorization_id"] == aid].iloc[0]
    if clean(req2["certificate_html"]):
        st.subheader("Certificate")
        st.html(req2["certificate_html"])
        st.download_button("Download Certificate", req2["certificate_html"], file_name=f"{req2['certificate_id']}.html", mime="text/html")


def crb_page(actor):
    st.header("Competency Review Board")
    auths = db_all("authorization_requests")
    pending = auths[auths["status"].isin(["QMS Reviewed","Technical Reviewed"])] if not auths.empty else pd.DataFrame()
    table(pending)
    if pending.empty:
        st.info("No pending CRB items.")
        return
    if actor_get(actor, "role") not in ["Admin","QMR","Technical Manager","Management","CRB Member","Tutor/Mentor"]:
        st.warning("Only CRB-related roles can submit CRB review.")
        return
    sel = st.selectbox("Review Request", pending["name"].astype(str)+" — "+pending["scope"].astype(str)+" — "+pending["authorization_id"].astype(str))
    aid = sel.split(" — ")[-1]
    req = pending[pending["authorization_id"] == aid].iloc[0]
    decision = st.selectbox("CRB Decision", ["Approved","Rejected","Deferred"])
    remarks = st.text_area("CRB Remarks")
    if st.button("Submit CRB Review"):
        db_insert("crb_reviews", {"crb_id": uid("CRB"), "authorization_id": aid, "user_id": req["user_id"], "name": req["name"], "scope": req["scope"], "review_date": today(), "tutor_decision": req["tutor_remarks"], "technical_decision": req["technical_remarks"], "qmr_decision": req["qms_remarks"], "management_decision": "", "final_decision": decision, "remarks": remarks, "signed_by": actor_get(actor,"name"), "created_on": now()})
        db_update("authorization_requests","authorization_id",aid,{"status":"CRB Approved" if decision=="Approved" else "CRB Rejected" if decision=="Rejected" else "CRB Deferred","crb_decision":decision,"crb_remarks":remarks,"updated_on":now()})
        st.success("CRB decision recorded.")
        st.rerun()
    st.subheader("CRB History")
    table(db_all("crb_reviews"))


def job_allocation_page(actor):
    st.header("Risk-Based Job Assignment Engine")
    with st.form("job"):
        title = st.text_input("Job Title")
        job_type = st.selectbox("Job Type", JOB_TYPES)
        scope = st.selectbox("Required Scope", SCOPES)
        vessel = st.text_input("Vessel / Project")
        imo = st.text_input("IMO Number")
        location = st.text_input("Location")
        planned = st.date_input("Planned Date")
        priority = st.selectbox("Priority", ["Low","Normal","High","Urgent"])
        risk = st.selectbox("Risk Level", ["Low","Medium","High","Critical"])
        min_level = st.selectbox("Minimum Competency Level", COMPETENCY_LEVELS, index=3)
        submit = st.form_submit_button("Create Job")
    if submit and title:
        db_insert("job_requests", {"job_id": uid("JOB"), "job_title": title, "job_type": job_type, "required_scope": scope, "vessel_name": vessel, "imo_number": imo, "location": location, "planned_date": str(planned), "priority": priority, "risk_level": risk, "minimum_level": min_level, "status": "Open", "created_by": actor_get(actor,"name"), "assigned_user_id": "", "assigned_user_name": "", "assignment_reason": "", "created_on": now(), "updated_on": now()})
        st.success("Job created.")
    jobs = db_all("job_requests")
    table(jobs)
    open_jobs = jobs[jobs["status"].isin(["Open","Reassign"])] if not jobs.empty else pd.DataFrame()
    if open_jobs.empty:
        return
    sel = st.selectbox("Select Job for Allocation", open_jobs["job_title"].astype(str)+" — "+open_jobs["job_id"].astype(str))
    jid = sel.split(" — ")[-1]
    job = jobs[jobs["job_id"] == jid].iloc[0]
    candidates = eligible_job_candidates(job)
    st.subheader("Eligible Candidates")
    table(candidates)
    if not candidates.empty:
        p = st.selectbox("Assign To", candidates["name"].astype(str)+" — "+candidates["user_id"].astype(str))
        uidv = p.split(" — ")[-1]
        cand = candidates[candidates["user_id"] == uidv].iloc[0]
        reason = f"Valid authorization {cand['authorization_id']}; scope {job['required_scope']}; level {cand['competency_level']}; KPI {cand['kpi_score']}; available."
        st.info(reason)
        if st.button("Assign Job"):
            db_update("job_requests","job_id",jid,{"status":"Assigned","assigned_user_id":uidv,"assigned_user_name":cand["name"],"assignment_reason":reason,"updated_on":now()})
            db_update("users","user_id",uidv,{"availability":"Busy"})
            st.success("Job assigned.")
    else:
        st.error("No eligible candidate found. Check authorization, scope, level, KPI, risk, and availability.")


def level_rank(level: str) -> int:
    m = re.search(r"Level\s+(\d+)", clean(level))
    return int(m.group(1)) if m else 0


def eligible_job_candidates(job: pd.Series) -> pd.DataFrame:
    auths = db_all("authorization_requests"); users = db_all("users"); kpis = db_all("kpi_records")
    if auths.empty or users.empty:
        return pd.DataFrame()
    approved = auths[(auths["status"] == "Management Approved") & (auths["scope"] == job["required_scope"]) & (auths["job_type"] == job["job_type"])]
    rows = []
    for _, a in approved.iterrows():
        if days_until(a["expiry_date"]) < 0:
            continue
        u = users[users["user_id"] == a["user_id"]]
        if u.empty:
            continue
        user = u.iloc[0]
        if user["availability"] != "Available":
            continue
        if level_rank(user["competency_level"]) < level_rank(job["minimum_level"]):
            continue
        user_kpis = kpis[kpis["user_id"] == user["user_id"]] if not kpis.empty else pd.DataFrame()
        kpi_score = float(user_kpis.sort_values("created_on").iloc[-1]["kpi_score"]) if not user_kpis.empty else 80.0
        if job["risk_level"] in ["High","Critical"] and kpi_score < 75:
            continue
        rows.append({"user_id": user["user_id"], "name": user["name"], "role": user["role"], "competency_level": user["competency_level"], "location": user["current_location"], "authorization_id": a["authorization_id"], "certificate_id": a["certificate_id"], "kpi_score": kpi_score})
    return pd.DataFrame(rows)


def kpi_page(actor):
    st.header("KPI and Utilization Engine")
    users = db_all("users")
    if actor_get(actor, "role") in ["Admin","Management","Technical Manager","QMR","Job Coordinator"]:
        with st.form("kpi"):
            person = st.selectbox("Person", users["name"].astype(str)+" — "+users["user_id"].astype(str)) if not users.empty else ""
            period = st.text_input("Period", datetime.now().strftime("%Y-%m"))
            surveys = st.number_input("Surveys Conducted", 0, 1000, 0)
            plans = st.number_input("Plans Reviewed", 0, 1000, 0)
            audits = st.number_input("Audits Done", 0, 1000, 0)
            overdue = st.number_input("Overdue Reports", 0, 1000, 0)
            ncr = st.number_input("NCR Count", 0, 1000, 0)
            feedback = st.slider("Client Feedback", 0.0, 100.0, 85.0)
            compliance = st.slider("Training Compliance", 0.0, 100.0, 90.0)
            utilization = st.slider("Utilization %", 0.0, 100.0, 70.0)
            remarks = st.text_area("Remarks")
            submit = st.form_submit_button("Save KPI")
        if submit and person:
            name, uidv = person.split(" — ")
            score = round((feedback*0.25 + compliance*0.25 + utilization*0.2 + max(0,100-overdue*5)*0.15 + max(0,100-ncr*10)*0.15), 2)
            db_insert("kpi_records", {"kpi_id": uid("KPI"), "user_id": uidv, "name": name, "period": period, "surveys_done": surveys, "plans_reviewed": plans, "audits_done": audits, "reports_overdue": overdue, "ncr_count": ncr, "client_feedback": feedback, "training_compliance": compliance, "utilization_percent": utilization, "kpi_score": score, "created_on": now(), "remarks": remarks})
            st.success(f"KPI saved. Score {score}")
    kpi = db_all("kpi_records")
    if actor_get(actor, "role") not in ["Admin","Management","Technical Manager","QMR","Job Coordinator"]:
        kpi = kpi[kpi["user_id"] == actor_get(actor, "user_id")] if not kpi.empty else kpi
    table(kpi)
    if not kpi.empty:
        st.bar_chart(kpi[["name","kpi_score"]].set_index("name"))


def cpd_page(actor):
    st.header("CPD / Seminars / Refresher Courses")
    users = db_all("users")
    with st.form("cpd"):
        if actor_get(actor, "role") in ["Admin","Trainer","QMR","Management"]:
            person = st.selectbox("Person", users["name"].astype(str)+" — "+users["user_id"].astype(str)) if not users.empty else ""
        else:
            person = f"{actor_get(actor,'name')} — {actor_get(actor,'user_id')}"
            st.text_input("Person", actor_get(actor,"name"), disabled=True)
        title = st.text_input("CPD / Seminar / Refresher Title")
        category = st.selectbox("Category", ["Seminar","Workshop","Webinar","Technical Update","Refresher Training","Conference"])
        hours = st.number_input("Hours", 0.0, 100.0, 2.0)
        provider = st.text_input("Provider", "PSB / BCS / External")
        completion = st.date_input("Completion Date")
        submit = st.form_submit_button("Add CPD")
    if submit and person and title:
        name, uidv = person.split(" — ")
        db_insert("cpd_records", {"cpd_id": uid("CPD"), "user_id": uidv, "name": name, "title": title, "category": category, "hours": hours, "provider": provider, "completion_date": str(completion), "evidence_file_id": "", "status": "Completed", "created_on": now()})
        st.success("CPD record added.")
    cpd = db_all("cpd_records")
    if actor_get(actor, "role") not in ["Admin","Trainer","QMR","Management"]:
        cpd = cpd[cpd["user_id"] == actor_get(actor,"user_id")] if not cpd.empty else cpd
    table(cpd)


def knowledge_page(actor):
    st.header("Technical Knowledge Library")
    if actor_get(actor, "role") in ["Admin","Trainer","QMR","Technical Manager","Rule Development Rep"]:
        with st.form("knowledge"):
            title = st.text_input("Title")
            category = st.selectbox("Category", ["Rule","Circular","Technical Bulletin","IMO Update","IACS Update","Interpretation","Lesson Learned"])
            standard = st.text_input("Standard / Reference")
            revision = st.text_input("Revision")
            mandatory = st.checkbox("Mandatory Acknowledgement", True)
            submit = st.form_submit_button("Add Knowledge Item")
        if submit and title:
            kid = uid("KNOW")
            db_insert("knowledge_library", {"knowledge_id": kid, "title": title, "category": category, "standard": standard, "revision": revision, "issue_date": today(), "file_id": "", "mandatory_ack": "Yes" if mandatory else "No", "uploaded_by": actor_get(actor,"name"), "created_on": now()})
            st.success("Knowledge item added. Upload file below if required.")
    lib = db_all("knowledge_library")
    table(lib)
    if not lib.empty:
        item = st.selectbox("Select Knowledge Item", lib["title"].astype(str)+" — "+lib["knowledge_id"].astype(str))
        kid = item.split(" — ")[-1]
        file_upload_panel(actor, "knowledge_library", kid, "Knowledge Bulletin")
        if st.button("Acknowledge Selected Item"):
            db_insert("knowledge_acknowledgements", {"ack_id": uid("ACK"), "knowledge_id": kid, "user_id": actor_get(actor,"user_id"), "name": actor_get(actor,"name"), "acknowledged_on": now(), "status": "Acknowledged"})
            st.success("Acknowledged.")
    st.subheader("Acknowledgements")
    table(db_all("knowledge_acknowledgements"))


def qms_page(actor):
    st.header("QMS / CAPA / Audit")
    tabs = st.tabs(["CAPA","Audit Trail","Notifications","Evidence Review"])
    with tabs[0]:
        users = db_all("users")
        with st.form("capa"):
            finding = st.text_input("Finding / NCR")
            severity = st.selectbox("Severity", ["Low","Medium","High","Critical"])
            owner = st.selectbox("Owner", users["name"].astype(str)+" — "+users["user_id"].astype(str)) if not users.empty else ""
            due = st.date_input("Due Date", date.today()+timedelta(days=30))
            action = st.text_area("Corrective Action")
            submit = st.form_submit_button("Create CAPA")
        if submit and finding and owner:
            name, uidv = owner.split(" — ")
            db_insert("capa_register", {"capa_id": uid("CAPA"), "source": "Training/Competency/QMS", "finding": finding, "severity": severity, "owner_id": uidv, "owner_name": name, "due_date": str(due), "status": "Open", "corrective_action": action, "created_on": now(), "updated_on": now()})
            st.success("CAPA created.")
        table(db_all("capa_register"))
    with tabs[1]:
        table(db_all("audit_trail"))
    with tabs[2]:
        table(db_all("notifications"))
    with tabs[3]:
        f = db_all("files")
        pending = f[f["review_status"] == "Pending Review"] if not f.empty else f
        table(pending)
        if not pending.empty:
            sel = st.selectbox("Review File", pending["file_name"].astype(str)+" — "+pending["file_id"].astype(str))
            fid = sel.split(" — ")[-1]
            status = st.selectbox("Review Status", ["Accepted","Rejected","Need Clarification"])
            if st.button("Save Review"):
                db_update("files","file_id",fid,{"review_status":status,"updated_on":now()})
                st.success("Review saved.")


def revalidation_page(actor):
    st.header("Revalidation / Reauthorization")
    auths = db_all("authorization_requests")
    approved = auths[auths["status"] == "Management Approved"] if not auths.empty else pd.DataFrame()
    if not approved.empty:
        approved = approved.copy()
        approved["days_to_expiry"] = approved["expiry_date"].apply(days_until)
        st.subheader("Expiring Authorizations")
        table(approved[approved["days_to_expiry"] <= 180])
        sel = st.selectbox("Select Authorization", approved["name"].astype(str)+" — "+approved["scope"].astype(str)+" — "+approved["authorization_id"].astype(str))
        aid = sel.split(" — ")[-1]
        req = approved[approved["authorization_id"] == aid].iloc[0]
        if st.button("Create Revalidation Request"):
            db_insert("revalidation_requests", {"revalidation_id": uid("REV"), "authorization_id": aid, "user_id": req["user_id"], "name": req["name"], "scope": req["scope"], "refresher_training_status": "Pending", "annual_review_status": "Pending", "kpi_review_status": "Pending", "tutor_confirmation": "Pending", "crb_status": "Pending", "final_status": "Open", "due_date": req["expiry_date"], "created_on": now(), "updated_on": now()})
            st.success("Revalidation request created.")
    table(db_all("revalidation_requests"))


def backup_page(actor):
    st.header("Audit Backup / Export")
    tables = ["users","training_modules","trainings","files","training_records","question_bank","assessment_history","competency_matrix","authorization_matrix","development_plans","field_exposure_matrix","witness_surveys","supervised_activities","authorization_requests","authorization_certificates","crb_reviews","annual_reviews","revalidation_requests","job_requests","kpi_records","cpd_records","knowledge_library","knowledge_acknowledgements","rule_library","document_versions","capa_register","notifications","audit_trail","technical_authorities","survey_report_reviews","plan_review_quality","competency_ncrs","authorization_restrictions","client_feedback","succession_plans","workforce_forecasts","accreditation_evidence","technical_interpretations"]
    export = {t: db_all(t).to_dict(orient="records") for t in tables}
    st.download_button("Download JSON Backup", json.dumps(export, indent=2, default=str), file_name=f"psb_hrdm_backup_{today()}.json", mime="application/json")
    with io.BytesIO() as buf:
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            for t in tables:
                db_all(t).to_excel(writer, sheet_name=t[:31], index=False)
        st.download_button("Download Excel Backup", buf.getvalue(), file_name=f"psb_hrdm_backup_{today()}.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


def qr_verify_page(actor):
    st.header("QR / Public Certificate Verification")
    cert_id = st.text_input("Certificate ID")
    if st.button("Verify") and cert_id:
        certs = db_all("authorization_certificates")
        c = certs[certs["certificate_id"] == cert_id] if not certs.empty else pd.DataFrame()
        if c.empty:
            st.error("Certificate not found.")
        else:
            row = c.iloc[0]
            if row["status"] == "Valid" and days_until(row["expiry_date"]) >= 0:
                st.success("Certificate is valid.")
            else:
                st.error("Certificate expired or invalid.")
            st.write(row[["certificate_id","name","job_type","scope","issue_date","expiry_date","status","verification_url"]])



def select_person(label, roles=None, key=None):
    users = db_all("users")
    if users.empty:
        return "", "", pd.Series(dtype=object)
    data = users if roles is None else users[users["role"].isin(roles)]
    if data.empty:
        return "", "", pd.Series(dtype=object)
    item = st.selectbox(label, data["name"].astype(str)+" — "+data["user_id"].astype(str), key=key)
    name, uidv = item.split(" — ")
    return name, uidv, data[data["user_id"] == uidv].iloc[0]

def technical_authority_page(actor):
    st.header("Technical Authority Framework")
    st.info("Register discipline technical authorities and approval limits.")
    if actor_get(actor,"role") in ["Admin","Management","Technical Manager"]:
        with st.form("ta"):
            name, uidv, _ = select_person("Authority Person", ["Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Technical Manager","Management"])
            discipline = st.selectbox("Discipline", ["Hull","Machinery","Electrical","Statutory","Plan Approval","Audit","Industrial","Rule Development"])
            level = st.selectbox("Authority Level", ["Discipline Expert","Principal","Head of Discipline","Technical Authority"])
            limit = st.text_area("Approval Limit", "Can approve technical interpretations and competency escalation within discipline.")
            remarks = st.text_area("Remarks")
            if st.form_submit_button("Appoint") and uidv:
                db_insert("technical_authorities", {"authority_id": uid("TA"), "user_id": uidv, "name": name, "discipline": discipline, "authority_level": level, "approval_limit": limit, "active": "Yes", "appointed_by": actor_get(actor,"name"), "appointed_on": today(), "remarks": remarks})
                st.success("Technical authority appointed.")
    table(db_all("technical_authorities"))

def survey_report_review_page(actor):
    st.header("Survey Report Review System")
    if actor_get(actor,"role") in ["Admin","Tutor/Mentor","Principal Surveyor","Technical Manager","QMR"]:
        with st.form("srr"):
            name, uidv, _ = select_person("Surveyor", ["Trainee","Surveyor","On Probation","Industrial Surveyor"])
            scope = st.selectbox("Survey Scope", SCOPES)
            vessel = st.text_input("Vessel / Project")
            file_id = st.text_input("Report File ID")
            tq = st.slider("Technical Quality",1,5,3); di=st.slider("Deficiency Identification",1,5,3)
            ri = st.slider("Rule Interpretation",1,5,3); rw=st.slider("Report Writing",1,5,3); dq=st.slider("Decision Quality",1,5,3)
            decision = st.selectbox("Decision", ["Accepted","Accepted with Comments","Rejected","Re-training Required"])
            comments = st.text_area("Comments")
            if st.form_submit_button("Save Review") and uidv:
                score = round((tq+di+ri+rw+dq)/25*100,2)
                db_insert("survey_report_reviews", {"review_id": uid("SRR"), "user_id": uidv, "name": name, "survey_scope": scope, "vessel_name": vessel, "report_file_id": file_id, "reviewer_id": actor_get(actor,"user_id"), "reviewer_name": actor_get(actor,"name"), "technical_quality": tq, "deficiency_identification": di, "rule_interpretation": ri, "report_writing": rw, "decision_quality": dq, "overall_score": score, "decision": decision, "comments": comments, "created_on": now()})
                if decision in ["Rejected","Re-training Required"]:
                    db_insert("competency_ncrs", {"ncr_id": uid("CNCR"), "user_id": uidv, "name": name, "source": "Survey Report Review", "scope": scope, "ncr_type": "Report Quality", "description": comments or decision, "severity": "Medium", "impact_on_authorization": "Review during revalidation", "status": "Open", "corrective_action": "Retraining/further supervision required", "raised_by": actor_get(actor,"name"), "raised_on": today(), "closed_on": ""})
                st.success(f"Review saved. Score {score}%")
    table(db_all("survey_report_reviews"))

def plan_review_quality_page(actor):
    st.header("Plan Review Quality Monitoring")
    if actor_get(actor,"role") in ["Admin","Chief Plan Appraiser","Technical Manager","QMR"]:
        with st.form("pqa"):
            name, uidv, _ = select_person("Plan Appraiser", ["Trainee","Plan Appraiser","On Probation"])
            scope = st.selectbox("Plan Scope", ["Plan Approval Hull","Plan Approval Machinery","Plan Approval Electrical"])
            project = st.text_input("Project / Drawing Package")
            file_id = st.text_input("Plan File ID")
            cq=st.slider("Comments Quality",1,5,3); missed=st.number_input("Missed Findings",0,100,0)
            turnaround=st.number_input("Turnaround Days",0,365,5); acc=st.slider("Accuracy",1,5,3)
            result=st.selectbox("Result", ["Accepted","Accepted with Comments","Rejected","Further Supervision Required"])
            comments=st.text_area("Comments")
            if st.form_submit_button("Save Plan QA") and uidv:
                score=max(0, round(((cq+acc)/10*100)-missed*5-max(0,turnaround-10),2))
                db_insert("plan_review_quality", {"planqa_id": uid("PQA"), "user_id": uidv, "name": name, "plan_scope": scope, "project_name": project, "plan_file_id": file_id, "reviewer_id": actor_get(actor,"user_id"), "reviewer_name": actor_get(actor,"name"), "comments_quality": cq, "missed_findings": missed, "turnaround_days": turnaround, "accuracy_score": acc, "overall_score": score, "result": result, "comments": comments, "created_on": now()})
                if result in ["Rejected","Further Supervision Required"] or missed>0:
                    db_insert("competency_ncrs", {"ncr_id": uid("CNCR"), "user_id": uidv, "name": name, "source": "Plan Review QA", "scope": scope, "ncr_type": "Plan Review Quality", "description": comments or result, "severity": "High" if missed>=3 else "Medium", "impact_on_authorization": "Affects revalidation/restriction", "status": "Open", "corrective_action": "Additional plan review supervision", "raised_by": actor_get(actor,"name"), "raised_on": today(), "closed_on": ""})
                st.success(f"Plan QA saved. Score {score}%")
    table(db_all("plan_review_quality"))

def competency_ncr_page(actor):
    st.header("Competency NCR / Performance Non-Conformance")
    if actor_get(actor,"role") in ["Admin","QMR","Technical Manager","Principal Surveyor","Chief Plan Appraiser","Lead Auditor"]:
        with st.form("cncr"):
            name, uidv, _ = select_person("Person", ["Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor","Trainee","On Probation"])
            source=st.selectbox("Source", ["Survey Decision","Missed Defect","Late Report","Rule Misinterpretation","Client Complaint","Audit Finding","Plan Review Error","Other"])
            scope=st.selectbox("Scope", SCOPES); severity=st.selectbox("Severity", ["Low","Medium","High","Critical"])
            impact=st.selectbox("Impact", ["None","Monitor","Restrict","Suspend","Withdraw","Re-training Required"])
            desc=st.text_area("Description"); action=st.text_area("Corrective Action")
            if st.form_submit_button("Raise NCR") and uidv:
                db_insert("competency_ncrs", {"ncr_id": uid("CNCR"), "user_id": uidv, "name": name, "source": source, "scope": scope, "ncr_type": source, "description": desc, "severity": severity, "impact_on_authorization": impact, "status": "Open", "corrective_action": action, "raised_by": actor_get(actor,"name"), "raised_on": today(), "closed_on": ""})
                st.success("Competency NCR raised.")
    ncrs=db_all("competency_ncrs"); table(ncrs)
    if actor_get(actor,"role") in ["Admin","QMR","Technical Manager"] and not ncrs.empty:
        open_ncr=ncrs[ncrs["status"]!="Closed"]
        if not open_ncr.empty:
            sel=st.selectbox("Close NCR", open_ncr["name"].astype(str)+" — "+open_ncr["ncr_id"].astype(str))
            if st.button("Close Selected NCR"):
                db_update("competency_ncrs","ncr_id",sel.split(" — ")[-1],{"status":"Closed","closed_on":today()}); st.rerun()

def competency_gap_advisor_page(actor):
    st.header("AI Competency Gap Advisor")
    name, uidv, _ = select_person("Person")
    scope=st.selectbox("Target Scope", SCOPES)
    if uidv:
        ok,gaps=readiness(uidv, scope)
        cpd=db_all("cpd_records"); ncrs=db_all("competency_ncrs")
        cpd_hours=cpd[cpd["user_id"]==uidv]["hours"].sum() if not cpd.empty else 0
        open_ncr=len(ncrs[(ncrs["user_id"]==uidv)&(ncrs["status"]!="Closed")]) if not ncrs.empty else 0
        if ok: st.success("No major authorization gap found.")
        else:
            st.error("Gaps found:")
            for g in gaps: st.write("- "+g)
        st.write(f"CPD hours: **{cpd_hours}** | Open competency NCRs: **{open_ncr}**")
        st.subheader("Recommended Actions")
        text=" ".join(gaps).lower()
        if "training" in text: st.write("- Complete/assign missing theoretical training modules.")
        if "witness" in text: st.write("- Tutor should schedule additional witness survey.")
        if "supervised" in text: st.write("- Tutor should schedule supervised survey/activity.")
        if "plan" in text: st.write("- Assign joint/independent plan review exercise.")
        if open_ncr: st.write("- Close competency NCRs before authorization or revalidation.")
        if cpd_hours < 20: st.write("- Complete annual CPD target, recommended minimum 20 hours.")

def annual_competency_board_page(actor):
    st.header("Annual Competency Review Board")
    if actor_get(actor,"role") in ["Admin","QMR","Technical Manager","Management","CRB Member"]:
        with st.form("ar"):
            name, uidv, _ = select_person("Person", ["Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor","Rule Development Rep","Trainee","On Probation"])
            scope=st.selectbox("Scope", SCOPES); year=st.number_input("Year",2020,2100,date.today().year)
            tr=st.selectbox("Training Status", ["Compliant","Partially Compliant","Non-Compliant"])
            kpi=st.selectbox("KPI Status", ["Good","Acceptable","Poor"]); comp=st.selectbox("Complaints", ["No Complaint","Minor","Major"])
            capa=st.selectbox("CAPA/NCR", ["No Open CAPA","Open Minor","Open Major"])
            decision=st.selectbox("Decision", ["Maintain","Upgrade","Restrict","Suspend","Withdraw","Additional Training"])
            remarks=st.text_area("Remarks")
            if st.form_submit_button("Save Annual Review") and uidv:
                db_insert("annual_reviews", {"review_id": uid("AR"), "user_id": uidv, "name": name, "scope": scope, "review_year": int(year), "training_status": tr, "kpi_status": kpi, "complaint_status": comp, "capa_status": capa, "decision": decision, "reviewer": actor_get(actor,"name"), "review_date": today(), "remarks": remarks})
                st.success("Annual review saved.")
    table(db_all("annual_reviews"))

def authorization_restrictions_page(actor):
    st.header("Authorization Restriction Matrix")
    auths=db_all("authorization_requests"); approved=auths[auths["status"]=="Management Approved"] if not auths.empty else pd.DataFrame()
    if actor_get(actor,"role") in ["Admin","QMR","Technical Manager","Management"] and not approved.empty:
        with st.form("res"):
            sel=st.selectbox("Authorization", approved["name"].astype(str)+" — "+approved["scope"].astype(str)+" — "+approved["authorization_id"].astype(str))
            aid=sel.split(" — ")[-1]; auth=approved[approved["authorization_id"]==aid].iloc[0]
            rtype=st.selectbox("Restriction Type", ["Scope Limit","Complexity Limit","Power/Capacity Limit","Only Under Supervision","Audit Type Limit","Temporary Restriction"])
            detail=st.text_area("Restriction Detail")
            eff=st.date_input("Effective Date"); exp=st.date_input("Expiry Date", date.today()+timedelta(days=365))
            if st.form_submit_button("Add Restriction"):
                db_insert("authorization_restrictions", {"restriction_id": uid("RES"), "authorization_id": auth["authorization_id"], "user_id": auth["user_id"], "name": auth["name"], "scope": auth["scope"], "restriction_type": rtype, "restriction_detail": detail, "effective_date": str(eff), "expiry_date": str(exp), "status": "Active", "imposed_by": actor_get(actor,"name"), "created_on": now()})
                st.success("Restriction added.")
    table(db_all("authorization_restrictions"))

def client_feedback_page(actor):
    st.header("Client / Shipowner / Shipyard Feedback")
    with st.form("fb"):
        name, uidv, _ = select_person("Person", ["Surveyor","Plan Appraiser","QMS Auditor","Industrial Surveyor"])
        client=st.text_input("Client / Shipowner / Shipyard"); project=st.text_input("Project / Vessel"); job_id=st.text_input("Job ID")
        rating=st.slider("Rating",1,5,4); ftype=st.selectbox("Type", ["Positive","Neutral","Complaint","Technical Concern"])
        comments=st.text_area("Comments"); impact=st.selectbox("KPI Impact", ["No Impact","Positive","Negative","Requires Review"])
        if st.form_submit_button("Save Feedback") and uidv:
            db_insert("client_feedback", {"feedback_id": uid("FB"), "user_id": uidv, "name": name, "client_name": client, "project_or_vessel": project, "job_id": job_id, "rating": rating, "feedback_type": ftype, "comments": comments, "impact_on_kpi": impact, "received_on": today()})
            if ftype in ["Complaint","Technical Concern"]:
                db_insert("competency_ncrs", {"ncr_id": uid("CNCR"), "user_id": uidv, "name": name, "source": "Client Feedback", "scope": "", "ncr_type": ftype, "description": comments, "severity": "Medium", "impact_on_authorization": "Review during annual competency review", "status": "Open", "corrective_action": "Investigate client feedback", "raised_by": actor_get(actor,"name"), "raised_on": today(), "closed_on": ""})
            st.success("Feedback saved.")
    table(db_all("client_feedback"))

def succession_planning_page(actor):
    st.header("Succession / Talent Pipeline")
    if actor_get(actor,"role") in ["Admin","Management","Technical Manager"]:
        with st.form("suc"):
            name, uidv, row = select_person("Person")
            target=st.selectbox("Target Role", ["Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Technical Manager","QMR","Management"])
            successor_for=st.text_input("Successor For / Position"); ready=st.selectbox("Readiness", ["Ready Now","Ready in 6 Months","Ready in 1 Year","Ready in 2 Years","Long-term Potential"])
            actions=st.text_area("Development Actions"); ready_date=st.date_input("Expected Ready Date", date.today()+timedelta(days=365)); sponsor=st.text_input("Sponsor", actor_get(actor,"name"))
            if st.form_submit_button("Save Succession Plan") and uidv:
                db_insert("succession_plans", {"succession_id": uid("SUC"), "user_id": uidv, "name": name, "current_role_name": row.get("role",""), "target_role": target, "readiness_level": ready, "successor_for": successor_for, "development_actions": actions, "expected_ready_date": str(ready_date), "sponsor": sponsor, "status": "Active", "created_on": now()})
                st.success("Succession plan saved.")
    table(db_all("succession_plans"))

def workforce_planning_page(actor):
    st.header("Workforce Planning / Resource Forecast")
    users=db_all("users"); auths=db_all("authorization_requests")
    if actor_get(actor,"role") in ["Admin","Management","Job Coordinator","Technical Manager"]:
        with st.form("wf"):
            period=st.text_input("Forecast Period", datetime.now().strftime("%Y-%m")); discipline=st.selectbox("Discipline", ["Hull","Machinery","Electrical","Statutory","Plan Approval","Audit","Industrial","Rule Development"])
            required=st.number_input("Required Headcount",0,1000,5); mitigation=st.text_area("Mitigation Plan")
            if st.form_submit_button("Save Forecast"):
                available=len(users[(users["availability"]=="Available")&(users["status"]=="Active")]) if not users.empty else 0
                expiring=sum(1 for _,a in auths.iterrows() if a.get("status","")=="Management Approved" and days_until(a.get("expiry_date",""))<=180) if not auths.empty else 0
                leave=len(users[users["availability"].isin(["On Leave","Unavailable"])]) if not users.empty else 0
                gap=int(required)-int(available); risk="High" if gap>0 or expiring>0 else "Low"
                db_insert("workforce_forecasts", {"forecast_id": uid("WF"), "forecast_period": period, "discipline": discipline, "required_headcount": int(required), "available_headcount": int(available), "expiring_authorizations": int(expiring), "leave_or_unavailable": int(leave), "gap": int(gap), "risk_status": risk, "mitigation_plan": mitigation, "created_on": now()})
                st.success("Forecast saved.")
    table(db_all("workforce_forecasts"))

def accreditation_readiness_page(actor):
    st.header("Accreditation Readiness Dashboard")
    if actor_get(actor,"role") in ["Admin","QMR","Management","Technical Manager"]:
        with st.form("acc"):
            standard=st.selectbox("Standard", ["IMO RO Code","ISO 9001","ISO/IEC 17020","IACS PR7","Internal QMS"])
            clause=st.text_input("Clause / Requirement Ref"); req=st.text_area("Requirement")
            linked_table=st.selectbox("Linked Evidence Table", ["training_records","competency_matrix","witness_surveys","supervised_activities","authorization_requests","authorization_certificates","kpi_records","cpd_records","capa_register","audit_trail","files"])
            linked_id=st.text_input("Linked Record ID"); summary=st.text_area("Evidence Summary"); owner=st.text_input("Owner", actor_get(actor,"name")); status=st.selectbox("Status", ["Ready","Partial","Gap","Not Applicable"])
            if st.form_submit_button("Save Evidence"):
                db_insert("accreditation_evidence", {"evidence_id": uid("ACC"), "standard": standard, "clause": clause, "requirement": req, "linked_table": linked_table, "linked_id": linked_id, "evidence_summary": summary, "status": status, "owner": owner, "last_reviewed": today()})
                st.success("Evidence saved.")
    evidence=db_all("accreditation_evidence"); table(evidence)
    if not evidence.empty:
        st.bar_chart(evidence.groupby("status").size().reset_index(name="count"), x="status", y="count")

def interpretation_portal_page(actor):
    st.header("Rule Interpretation / Technical Decision Portal")
    if actor_get(actor,"role") in ["Admin","Technical Manager","Principal Surveyor","Chief Plan Appraiser","Lead Auditor","Rule Development Rep"]:
        with st.form("interp"):
            title=st.text_input("Title"); discipline=st.selectbox("Discipline", ["Hull","Machinery","Electrical","Statutory","Plan Approval","Audit","Industrial","Rule Development"])
            related=st.text_input("Related Rule / Clause"); question=st.text_area("Question / Case"); interpretation=st.text_area("Approved Interpretation / Decision")
            approved=st.text_input("Approved By", actor_get(actor,"name")); status=st.selectbox("Status", ["Draft","Approved","Withdrawn","Superseded"]); rev=st.text_input("Revision", "Rev.0"); issue=st.date_input("Issue Date")
            if st.form_submit_button("Save Interpretation") and title:
                db_insert("technical_interpretations", {"interpretation_id": uid("INT"), "title": title, "discipline": discipline, "related_rule": related, "question": question, "interpretation": interpretation, "approved_by": approved, "approval_status": status, "revision": rev, "issue_date": str(issue), "created_on": now()})
                st.success("Interpretation saved.")
    table(db_all("technical_interpretations"))

def management_page(actor):
    st.header("Management Dashboard")
    dashboard_page(actor)
    st.subheader("Authorizations")
    table(db_all("authorization_requests"))
    st.subheader("Jobs")
    table(db_all("job_requests"))
    st.subheader("KPI")
    kpi = db_all("kpi_records")
    table(kpi)
    if not kpi.empty:
        st.bar_chart(kpi[["name","kpi_score"]].set_index("name"))


def main() -> None:
    st.set_page_config(page_title=APP_TITLE, page_icon="⚓", layout="wide", initial_sidebar_state="expanded")
    apply_style()
    require_persistent_backend()
    init_db()
    actor = require_login()
    header()
    page = sidebar(actor)
    if page == "Dashboard": dashboard_page(actor)
    elif page == "Admin": admin_page(actor)
    elif page == "Training Matrix": training_matrix_page(actor)
    elif page == "Training": training_page(actor)
    elif page == "Files": files_page(actor)
    elif page == "Development Plans": development_plan_page(actor)
    elif page == "Competency": competency_page(actor)
    elif page == "Practical/Witness": practical_page(actor)
    elif page == "Authorization": authorization_page(actor)
    elif page == "CRB": crb_page(actor)
    elif page == "Job Allocation": job_allocation_page(actor)
    elif page == "KPI": kpi_page(actor)
    elif page == "CPD": cpd_page(actor)
    elif page == "Knowledge Library": knowledge_page(actor)
    elif page == "QMS": qms_page(actor)
    elif page == "Revalidation": revalidation_page(actor)
    elif page == "Backup": backup_page(actor)
    elif page == "QR Verify": qr_verify_page(actor)
    elif page == "Technical Authority": technical_authority_page(actor)
    elif page == "Survey Report Review": survey_report_review_page(actor)
    elif page == "Plan Review QA": plan_review_quality_page(actor)
    elif page == "Competency NCR": competency_ncr_page(actor)
    elif page == "Gap Advisor": competency_gap_advisor_page(actor)
    elif page == "Annual Board": annual_competency_board_page(actor)
    elif page == "Restrictions": authorization_restrictions_page(actor)
    elif page == "Client Feedback": client_feedback_page(actor)
    elif page == "Succession": succession_planning_page(actor)
    elif page == "Workforce Planning": workforce_planning_page(actor)
    elif page == "Accreditation Readiness": accreditation_readiness_page(actor)
    elif page == "Interpretation Portal": interpretation_portal_page(actor)
    elif page == "Management": management_page(actor)


if __name__ == "__main__":
    main()
